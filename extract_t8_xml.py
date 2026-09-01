from pptx import Presentation
import xml.etree.ElementTree as ET

prs = Presentation("data/output/t8.pptx")
with open("t8_shapes.xml", "w", encoding="utf-8") as f:
    s1 = prs.slides[0]
    for sh in s1.shapes:
        if sh.name == "Rectangle 14":
            f.write(ET.tostring(sh._element, encoding="utf-8").decode("utf-8"))
print("Saved t8_shapes.xml")
