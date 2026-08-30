"""
render_cache.py
===============
Tile and shape raster memory cache & multi-slide parallel rendering worker pool.
"""

from __future__ import annotations

import os
import io
import logging
import threading
from collections import OrderedDict
from typing import Dict, List, Optional, Tuple, Any
from concurrent.futures import ThreadPoolExecutor
from PIL import Image
from pptx import Presentation

log = logging.getLogger("pptx_renderers.cache")


class RenderCache:
    """Thread-safe bounded memory cache for theme palettes, fonts, and rasterized shape layers."""

    _lock = threading.Lock()
    _max_entries = 128
    _theme_cache: OrderedDict[str, Any] = OrderedDict()
    _font_cache: OrderedDict[str, Any] = OrderedDict()
    _shape_cache: OrderedDict[str, Image.Image] = OrderedDict()

    @classmethod
    def get_theme(cls, key: str) -> Optional[Any]:
        with cls._lock:
            if key in cls._theme_cache:
                cls._theme_cache.move_to_end(key)
                return cls._theme_cache[key]
            return None

    @classmethod
    def set_theme(cls, key: str, theme: Any) -> None:
        with cls._lock:
            cls._theme_cache[key] = theme
            cls._theme_cache.move_to_end(key)
            if len(cls._theme_cache) > cls._max_entries:
                cls._theme_cache.popitem(last=False)

    @classmethod
    def clear(cls) -> None:
        with cls._lock:
            cls._theme_cache.clear()
            cls._font_cache.clear()
            cls._shape_cache.clear()


def render_single_slide_task(slide_idx: int, pptx_source: Any, width: int) -> Tuple[int, Image.Image]:
    """Helper worker task for parallel slide rendering."""
    from pptx_jahat.tools.preview import render_slide, Theme, FontResolver, _theme_for_slide
    prs = Presentation(pptx_source)
    slide = prs.slides[slide_idx]
    theme = _theme_for_slide(slide, {}, Theme.from_presentation(prs))
    fonts = FontResolver()
    fonts.set_theme_fonts(theme.major_font, theme.minor_font)
    img = render_slide(slide, prs, width=width, theme=theme, fonts=fonts)
    return slide_idx, img


def render_pptx_parallel(source: Any, width: int = 1280,
                         slide_numbers: Optional[List[int]] = None,
                         max_workers: Optional[int] = None) -> List[Image.Image]:
    """
    Renders multiple slides concurrently across a multi-core worker pool.
    """
    # Read into memory buffer so threads can safely open their own Presentation instances
    if isinstance(source, (str, os.PathLike)):
        with open(source, "rb") as f:
            pptx_bytes = f.read()
    elif isinstance(source, io.BytesIO):
        pptx_bytes = source.getvalue()
    elif isinstance(source, bytes):
        pptx_bytes = source
    else:
        # Fallback to sequential
        from pptx_jahat.tools.preview import render_pptx
        return render_pptx(source, width=width, slide_numbers=slide_numbers)

    prs = Presentation(io.BytesIO(pptx_bytes))
    total_slides = len(prs.slides)
    target_indices = [
        i for i in range(total_slides)
        if slide_numbers is None or (i + 1) in slide_numbers
    ]

    if not target_indices:
        return []

    workers = max_workers or min(4, (os.cpu_count() or 1))
    results: List[Tuple[int, Image.Image]] = []

    with ThreadPoolExecutor(max_workers=workers) as executor:
        futures = [
            executor.submit(render_single_slide_task, idx, io.BytesIO(pptx_bytes), width)
            for idx in target_indices
        ]
        for fut in futures:
            results.append(fut.result())

    # Sort results by original slide order
    results.sort(key=lambda r: r[0])
    return [img for _, img in results]
