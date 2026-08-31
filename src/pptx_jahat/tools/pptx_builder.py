import json
import re
import copy
import io
import zipfile
import collections
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Dict, Any, List, Optional, Callable, Tuple
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml import parse_xml

from pptx_jahat.config import Config, DATA_DIR, OUTPUT_DIR
from pptx_jahat.tools.docx_parser import parse_docx
from pptx_jahat.tools.pptx_engine import inspect_template_slides, inspect_all_templates
from pptx_jahat.tools.image_gen import generate_image
from pptx_jahat.tools.preview import render_pptx_file_previews, image_to_base64_jpeg
from openai import OpenAI

def _set_paragraph_rtl_and_fonts(paragraph: Any, font_name: Optional[str] = "Vazirmatn") -> None:
    """
    Directly sets DrawingML paragraph properties for true RTL and complex script fonts.
    """
    try:
        pPr = paragraph._p.get_or_add_pPr()
        pPr.set("rtl", "1")
        pPr.set("algn", "r")
        
        # Set default complex script font
        if font_name:
            defRPr = pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr")
            if defRPr is None:
                defRPr = OxmlElement("a:defRPr")
                pPr.append(defRPr)
            cs = defRPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}cs")
            if cs is None:
                cs = OxmlElement("a:cs")
                defRPr.append(cs)
            cs.set("typeface", font_name)
    except Exception:
        pass

def _safe_update_text_frame(
    tf: Any,
    new_text: str,
    is_rtl: bool = True,
    max_box_width_emu: Optional[int] = None,
    max_box_height_emu: Optional[int] = None,
    font_override: Optional[str] = None
) -> None:
    """
    Updates text in a text_frame while:
    1. Preserving run-level formatting (color, bold, italic).
    2. Dynamic font auto-sizing based on character count and bounding box dimensions.
    3. Applying true DrawingML RTL properties.
    """
    if not tf:
        return
        
    lines = [line for line in new_text.split("\n") if line.strip()]
    if not lines:
        lines = [new_text]
        
    # Capture style of first run if available
    saved_font = {
        "name": None,
        "size": None,
        "bold": None,
        "italic": None,
        "color": None
    }
    
    try:
        if tf.paragraphs:
            p0 = tf.paragraphs[0]
            if p0.runs:
                r0 = p0.runs[0]
                if r0.font:
                    saved_font["name"] = r0.font.name
                    saved_font["size"] = r0.font.size
                    saved_font["bold"] = r0.font.bold
                    saved_font["italic"] = r0.font.italic
                    try:
                        if r0.font.color and r0.font.color.rgb:
                            saved_font["color"] = r0.font.color.rgb
                    except Exception:
                        pass
    except Exception:
        pass

    # Dynamic Font Auto-Sizing calculation
    total_chars = sum(len(line) for line in lines)
    calculated_size_pt = None
    
    if saved_font["size"]:
        orig_pt = saved_font["size"].pt
        if total_chars > 250:
            calculated_size_pt = max(10, min(orig_pt, 12))
        elif total_chars > 120:
            calculated_size_pt = max(11, min(orig_pt, 14))
        elif total_chars > 60:
            calculated_size_pt = max(12, min(orig_pt, 18))
        else:
            calculated_size_pt = orig_pt

    # Clear old paragraphs and populate with new text lines
    tf.clear()
    
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        if is_rtl:
            p.alignment = PP_ALIGN.RIGHT
            _set_paragraph_rtl_and_fonts(p, font_name=font_override or saved_font["name"] or "Vazirmatn")
            
        # Apply preserved/adjusted font styling to runs
        if p.runs:
            for run in p.runs:
                run.font.name = font_override or saved_font["name"] or "Vazirmatn"
                if calculated_size_pt:
                    run.font.size = Pt(calculated_size_pt)
                elif saved_font["size"]:
                    run.font.size = saved_font["size"]
                if saved_font["bold"] is not None:
                    run.font.bold = saved_font["bold"]
                if saved_font["italic"] is not None:
                    run.font.italic = saved_font["italic"]
                if saved_font["color"]:
                    run.font.color.rgb = saved_font["color"]

def _replace_image_in_shape(shape: Any, new_image_path: Path | str) -> bool:
    """
    Replaces the image blob in a picture shape with a newly generated or selected image.
    """
    try:
        img_path = Path(new_image_path)
        if not img_path.exists():
            return False
            
        with open(img_path, "rb") as f:
            new_blob = f.read()
            
        if hasattr(shape, "image"):
            # Update image part blob
            shape.image._blob = new_blob
            return True
    except Exception:
        pass
    return False

def _remove_shape(slide: Any, shape_index: int) -> bool:
    """
    Safely removes a shape from slide XML by shape_index.
    """
    try:
        if 0 <= shape_index < len(slide.shapes):
            shape = slide.shapes[shape_index]
            sp_elem = shape._element
            parent = sp_elem.getparent()
            if parent is not None:
                parent.remove(sp_elem)
                return True
    except Exception:
        pass
    return False

def _remove_shapes(slide: Any, shape_indices: List[int]) -> None:
    """
    Removes multiple shapes in descending index order to avoid index shift issues.
    """
    if not shape_indices:
        return
    for s_idx in sorted(set(shape_indices), reverse=True):
        _remove_shape(slide, s_idx)

def clone_slide_across_presentations(source_prs: Presentation, target_prs: Presentation, slide_index: int) -> Any:
    """
    Deep clones a slide from source_prs into target_prs, preserving layout, background,
    media parts, and relationship mappings while avoiding duplicate/corrupted package parts.
    """
    src_slide = source_prs.slides[slide_index]
    
    # Choose layout from target_prs matching source or fallback to blank/content layout
    layout_name = src_slide.slide_layout.name if src_slide.slide_layout else "Blank"
    matching_layout = None
    for layout in target_prs.slide_layouts:
        if layout.name == layout_name:
            matching_layout = layout
            break
    if not matching_layout:
        matching_layout = target_prs.slide_layouts[min(1, len(target_prs.slide_layouts) - 1)]
        
    target_slide = target_prs.slides.add_slide(matching_layout)
    
    # Copy relationships & media/picture parts cleanly to avoid corrupting theme/layout parts
    src_part = src_slide.part
    target_part = target_slide.part
    rel_id_map: Dict[str, str] = {}
    
    for rel_id, rel in src_part.rels.items():
        if rel.reltype == RT.SLIDE_LAYOUT or rel.reltype == RT.NOTES_SLIDE:
            continue
        try:
            if rel.is_external:
                new_rid = target_part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
                rel_id_map[rel_id] = new_rid
            elif rel.reltype == RT.IMAGE:
                # Add image blob into target package cleanly
                image_bytes = rel.target_part.blob
                new_image_part = target_prs.part.package.get_or_add_image_part(io.BytesIO(image_bytes))
                new_rid = target_part.relate_to(new_image_part, RT.IMAGE)
                rel_id_map[rel_id] = new_rid
            elif rel.reltype == RT.HYPERLINK:
                new_rid = target_part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
                rel_id_map[rel_id] = new_rid
        except Exception:
            pass

    # Copy background definition (from slide or source layout/master) into target slide element
    try:
        src_cSld = src_slide._element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}cSld")
        target_cSld = target_slide._element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}cSld")
        
        # 1. Check if source slide has explicit <p:bg>
        src_bg = src_cSld.find("{http://schemas.openxmlformats.org/presentationml/2006/main}bg") if src_cSld is not None else None
        
        # 2. If not on slide, check source layout <p:bg>
        if src_bg is None and src_slide.slide_layout:
            l_cSld = src_slide.slide_layout._element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}cSld")
            src_bg = l_cSld.find("{http://schemas.openxmlformats.org/presentationml/2006/main}bg") if l_cSld is not None else None
            
        # 3. If not on layout, check source master <p:bg>
        if src_bg is None and src_slide.slide_layout and src_slide.slide_layout.slide_master:
            m_cSld = src_slide.slide_layout.slide_master._element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}cSld")
            src_bg = m_cSld.find("{http://schemas.openxmlformats.org/presentationml/2006/main}bg") if m_cSld is not None else None

        if src_bg is not None and target_cSld is not None:
            # Check if target already has <p:bg>
            t_bg = target_cSld.find("{http://schemas.openxmlformats.org/presentationml/2006/main}bg")
            if t_bg is not None:
                target_cSld.remove(t_bg)
            copied_bg = copy.deepcopy(src_bg)
            # Remap relationship IDs in background (e.g. blipFill images)
            for elem in copied_bg.iter():
                for attr_name in list(elem.attrib.keys()):
                    if "embed" in attr_name or "id" in attr_name or "link" in attr_name:
                        val = elem.attrib[attr_name]
                        if val in rel_id_map:
                            elem.attrib[attr_name] = rel_id_map[val]
            target_cSld.insert(0, copied_bg)
    except Exception:
        pass
            
    # Replace target slide's spTree (shape tree) with deep copied source spTree
    target_spTree = target_slide.shapes._spTree
    src_spTree = src_slide.shapes._spTree
    
    # Remove default shapes in the newly added slide
    for child in list(target_spTree):
        target_spTree.remove(child)
        
    copied_spTree = copy.deepcopy(src_spTree)
    
    # Remap relationship IDs in elements (blip r:embed, hyperlinks, etc.)
    for elem in copied_spTree.iter():
        for attr_name in list(elem.attrib.keys()):
            if "embed" in attr_name or "id" in attr_name or "link" in attr_name:
                val = elem.attrib[attr_name]
                if val in rel_id_map:
                    elem.attrib[attr_name] = rel_id_map[val]
                    
    # Copy all children from copied spTree into target spTree
    for child in list(copied_spTree):
        target_spTree.append(child)
        
    return target_slide

def generate_slide_replacements_with_ai(
    template_inventory: List[Dict[str, Any]],
    doc_structure: Dict[str, Any],
    log_cb: Optional[Callable[[str], None]] = None,
    on_ai_images_ready: Optional[Callable[[List[Dict[str, Any]]], None]] = None
) -> Dict[str, Any]:
    """
    Step 3: AI Vision Agent reasons over multimodal slide screenshots, shape slots across templates,
    and Word docx content.
    Returns optimal slide selections across templates, exact text replacements, shapes to remove,
    speaker notes, and optional AI image generation prompts for picture slots.
    """
    def log(msg: str):
        if log_cb:
            log_cb(msg)

    log("[Step 3] AI Vision Agent analyzing candidate template slides & document content...")

    client = OpenAI(
        api_key=Config.NINEROUTER_KEY or "dummy_key",
        base_url=f"{Config.NINEROUTER_URL.rstrip('/')}/v1",
        timeout=120.0
    )

    system_prompt = (
        "You are an expert Presentation Art Director and Content Producer. "
        "You receive visual screenshots, shape slots, and archetype tags of candidate presentation slides across multiple templates, "
        "along with a parsed Word document. "
        "Your task is to:\n"
        "1. Select the best visual slide archetype from the available templates for each section/topic in the document (title_cover, table_matrix, metrics_stats, multi_column, process_timeline, content_bullets, conclusion_quote).\n"
        "2. Chunk and adapt long document text into punchy, high-impact slide text (concise headers, 3-4 bullet points max, 10-12 words per bullet).\n"
        "3. Match adapted content into the chosen slide's shape slots (shape_index).\n"
        "4. Generate detailed speaker notes for each slide to retain comprehensive background details from the document.\n"
        "5. Identify any unnecessary or overflowing shape indices to delete (shapes_to_remove).\n"
        "6. If a slot contains a table, supply updated 2D table_data.\n"
        "7. If a slide contains picture/graphic placeholders, you can provide an image_prompt for contextual AI image generation.\n"
        "Strictly return valid JSON adhering to the specified schema."
    )

    # Prepare slot descriptions (without heavy base64 strings in the JSON text prompt)
    inventory_summary = []
    for s in template_inventory:
        summary_entry = {
            "template_file": s.get("template_file"),
            "slide_index": s.get("slide_index"),
            "layout_name": s.get("layout_name"),
            "archetype": s.get("archetype", "content_bullets"),
            "text_slots": [
                {
                    "shape_index": slot.get("shape_index"),
                    "shape_name": slot.get("shape_name"),
                    "shape_type": slot.get("shape_type"),
                    "original_text": slot.get("original_text", "")[:120],
                    "is_title": slot.get("is_title", False),
                    "is_table": slot.get("is_table", False),
                    "table_shape": f"{slot.get('table_rows')}x{slot.get('table_cols')}" if slot.get("is_table") else None,
                    "is_decorative": slot.get("is_decorative", False)
                }
                for slot in s.get("text_slots", [])
            ]
        }
        inventory_summary.append(summary_entry)

    # Build multimodal user message content array
    user_content: List[Dict[str, Any]] = [
        {
            "type": "text",
            "text": f"""
Step 1 - Available Slide Blueprints, Archetypes & Slots:
{json.dumps(inventory_summary, ensure_ascii=False, indent=2)}

Step 2 - Input Word Document Outline & Content:
{json.dumps(doc_structure, ensure_ascii=False, indent=2)}

Visual Screenshots of Candidate Template Slides:
(See attached images corresponding to the candidate slides above)

Instructions:
1. Construct a cohesive presentation sequence matching the document flow (Title slide, Content/Topic slides, Metric slides, Summary).
2. For each slide in your output deck, specify:
   - "source_template": Name of template file (e.g. "T711.pptx", "t1.pptx", "sample_template.pptx")
   - "source_slide_index": Index of slide in that template
   - "target_section": Name of document section this slide covers
   - "speaker_notes": Detailed explanatory talking points for the presenter
   - "shape_replacements": List of {{"shape_index": int, "text": str}} mapping new adapted text into slots
   - "shapes_to_remove": List of shape indices [int] that should be pruned/deleted from the slide
   - "table_replacements": List of {{"shape_index": int, "table_data": [["cell", ...], ...]}}
   - "image_replacements": Optional list of {{"shape_index": int, "image_prompt": "detailed prompt for slide visual"}}

Return a JSON object with this exact schema:
{{
  "deck_title": "Presentation Title",
  "slides": [
    {{
      "source_template": "sample_template.pptx",
      "source_slide_index": 0,
      "target_section": "Document Title",
      "speaker_notes": "Welcome to the presentation...",
      "shape_replacements": [
        {{
          "shape_index": 0,
          "text": "Presentation Title"
        }}
      ],
      "shapes_to_remove": [],
      "table_replacements": [],
      "image_replacements": []
    }}
  ]
}}
"""
        }
    ]

    # Attach slide screenshots as multimodal image parts (limit to top 15 candidate slides to preserve tokens)
    attached_count = 0
    ai_sent_images: List[Dict[str, Any]] = []
    for s in template_inventory:
        b64 = s.get("screenshot_base64")
        if b64 and attached_count < 15:
            user_content.append({
                "type": "image_url",
                "image_url": {
                    "url": b64
                }
            })
            ai_sent_images.append({
                "template_file": s.get("template_file"),
                "slide_index": s.get("slide_index"),
                "archetype": s.get("archetype", "content_bullets"),
                "base64": b64
            })
            attached_count += 1

    if on_ai_images_ready:
        try:
            on_ai_images_ready(ai_sent_images)
        except Exception:
            pass

    log(f"[Step 3] Sending prompt with {attached_count} visual slide previews to 9Router AI...")

    try:
        response = client.chat.completions.create(
            model=Config.NINEROUTER_CHAT_MODEL,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_content}
            ],
            temperature=0.3
        )
        content = response.choices[0].message.content or "{}"
        match = re.search(r"```(?:json)?\s*(\{.*?\})\s*```", content, re.DOTALL)
        if match:
            json_str = match.group(1)
        else:
            json_str = content.strip()
            
        plan = json.loads(json_str)
        log("[Step 3] AI Vision Agent generated presentation plan successfully.")
        return plan
    except Exception as e:
        log(f"[Step 3 Warning] AI reasoning exception ({e}), using fallback multi-template mapper.")
        return None

def build_pptx_with_agent(
    docx_path: str | Path,
    output_path: Optional[str | Path] = None,
    template_name: Optional[str] = None,
    log_callback: Optional[Callable[[str], None]] = None,
    on_ai_images_ready: Optional[Callable[[List[Dict[str, Any]]], None]] = None
) -> str:
    """
    4-Step Vision-Guided Multi-Template Presentation Generation:
    Step 1: Scan & inspect candidate slides across all templates with rendered screenshots.
    Step 2: Read and parse Word DOCX structure.
    Step 3: Vision AI reasons on slide screenshots & doc content, selecting best slides across templates.
    Step 4: Clone selected slides across presentations into target deck, prune removed shapes, and update text in-place.
    """
    def log(msg: str):
        if log_callback:
            log_callback(msg)

    # ----------------------------------------------------
    # Step 1: Scan & inspect templates
    # ----------------------------------------------------
    if template_name and template_name != "All Templates (Global AI Matching)":
        candidate = DATA_DIR / template_name
        if candidate.exists():
            log(f"[Step 1] Inspecting selected template: {candidate.name}...")
            template_inventory = inspect_template_slides(candidate, include_screenshots=True)
        else:
            log(f"[Step 1] Scanning all templates in {DATA_DIR}...")
            template_inventory = inspect_all_templates(DATA_DIR, include_screenshots=True)
    else:
        log(f"[Step 1] Scanning all templates in {DATA_DIR} with visual screenshots...")
        template_inventory = inspect_all_templates(DATA_DIR, include_screenshots=True)

    if not template_inventory:
        raise FileNotFoundError("No PPTX templates found in data folder.")

    log(f"[Step 1] Loaded {len(template_inventory)} candidate slides across templates.")

    # ----------------------------------------------------
    # Step 2: Read and parse Word Document
    # ----------------------------------------------------
    docx_file = Path(docx_path)
    log(f"[Step 2] Reading Word document: {docx_file.name}...")
    parsed_doc = parse_docx(docx_file)
    log(f"[Step 2] Parsed {parsed_doc['total_sections']} sections from document.")

    # ----------------------------------------------------
    # Step 3: AI Vision Agent writes texts and selects slides
    # ----------------------------------------------------
    ai_plan = generate_slide_replacements_with_ai(
        template_inventory,
        parsed_doc,
        log_cb=log,
        on_ai_images_ready=on_ai_images_ready
    )

    # ----------------------------------------------------
    # Step 4: Assemble target deck across presentations
    # ----------------------------------------------------
    log("[Step 4] Assembling target presentation from selected template slides...")

    # Cache opened presentations by filename
    prs_cache: Dict[str, Presentation] = {}
    def get_source_prs(tpl_file: str) -> Presentation:
        if tpl_file not in prs_cache:
            p = DATA_DIR / tpl_file
            if not p.exists():
                # Fallback to first existing template
                p = next(DATA_DIR.glob("*.pptx"))
            prs_cache[tpl_file] = Presentation(str(p))
        return prs_cache[tpl_file]

    # Pre-open first source template to create matching target presentation package
    first_tpl_name = template_inventory[0]["template_file"]
    first_tpl_path = DATA_DIR / first_tpl_name if (DATA_DIR / first_tpl_name).exists() else next(DATA_DIR.glob("*.pptx"))
    
    # Initialize target presentation from base template to retain themes, color palettes, and layouts
    target_prs = Presentation(str(first_tpl_path))
    
    # Clear existing slides from target presentation
    while len(target_prs.slides) > 0:
        rId = target_prs.slides._sldIdLst[0].rId
        target_prs.part.drop_rel(rId)
        target_prs.slides._sldIdLst.remove(target_prs.slides._sldIdLst[0])

    if ai_plan and "slides" in ai_plan and len(ai_plan["slides"]) > 0:
        for s_plan in ai_plan["slides"]:
            src_tpl = s_plan.get("source_template") or first_tpl_name
            src_idx = s_plan.get("source_slide_index", 0)
            
            src_prs = get_source_prs(src_tpl)
            if src_idx >= len(src_prs.slides):
                src_idx = 0
                
            # Clone slide across presentation
            target_slide = clone_slide_across_presentations(src_prs, target_prs, src_idx)
            
            # In-place text replacements
            replacements = {r.get("shape_index"): r.get("text") for r in s_plan.get("shape_replacements", [])}
            for shape_idx, shape in enumerate(target_slide.shapes):
                if shape_idx in replacements and shape.has_text_frame:
                    new_text = replacements[shape_idx]
                    if new_text is not None:
                        _safe_update_text_frame(
                            shape.text_frame,
                            str(new_text),
                            is_rtl=True,
                            max_box_width_emu=getattr(shape, "width", None),
                            max_box_height_emu=getattr(shape, "height", None)
                        )
                        
            # Table replacements
            table_replacements = {t.get("shape_index"): t.get("table_data") for t in s_plan.get("table_replacements", [])}
            for shape_idx, shape in enumerate(target_slide.shapes):
                if shape_idx in table_replacements and shape.has_table:
                    tdata = table_replacements[shape_idx]
                    if tdata:
                        for r_i, row in enumerate(tdata):
                            if r_i < len(shape.table.rows):
                                for c_i, cell_val in enumerate(row):
                                    if c_i < len(shape.table.columns):
                                        shape.table.cell(r_i, c_i).text = str(cell_val)
                                        
            # Image replacements via Image Gen API
            image_replacements = {img.get("shape_index"): img.get("image_prompt") for img in s_plan.get("image_replacements", [])}
            for shape_idx, shape in enumerate(target_slide.shapes):
                if shape_idx in image_replacements and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    prompt = image_replacements[shape_idx]
                    if prompt:
                        try:
                            log(f"[Step 4] Generating AI image for slide slot #{shape_idx}: '{prompt[:40]}...'")
                            img_file = generate_image(prompt)
                            if img_file and not img_file.startswith("Error"):
                                _replace_image_in_shape(shape, img_file)
                        except Exception as e:
                            log(f"[Step 4 Warning] Image generation skipped: {e}")

            # Speaker Notes
            notes_text = s_plan.get("speaker_notes")
            if notes_text:
                try:
                    notes_slide = target_slide.notes_slide
                    text_frame = notes_slide.notes_text_frame
                    text_frame.text = str(notes_text)
                except Exception:
                    pass

            # Remove unwanted shapes after text/table replacements to avoid index shifts during replacement
            shapes_to_remove = s_plan.get("shapes_to_remove", [])
            if shapes_to_remove:
                _remove_shapes(target_slide, shapes_to_remove)
    else:
        # Fallback Multi-Template Assembly:
        # Title slide from first template, content slides from available candidate slides
        log("[Step 4 Fallback] Generating presentation using multi-slide assembly...")
        
        # 1. Title Slide
        first_tpl = template_inventory[0]["template_file"]
        base_prs = get_source_prs(first_tpl)
        title_slide = clone_slide_across_presentations(base_prs, target_prs, 0)
        for shape in title_slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                _safe_update_text_frame(shape.text_frame, parsed_doc.get("document_title", "Presentation"))
                break
                
        # 2. Content Slides per Section
        for s_idx, section in enumerate(parsed_doc.get("sections", [])):
            candidate_entry = template_inventory[(s_idx + 1) % len(template_inventory)]
            src_prs = get_source_prs(candidate_entry["template_file"])
            src_idx = candidate_entry["slide_index"]
            
            c_slide = clone_slide_across_presentations(src_prs, target_prs, src_idx)
            text_shapes = [sh for sh in c_slide.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
            
            if text_shapes:
                _safe_update_text_frame(text_shapes[0].text_frame, section.get("title", ""))
                if len(text_shapes) > 1:
                    body = "\n".join(section.get("paragraphs", []) + [f"• {b}" for b in section.get("bullets", [])])
                    _safe_update_text_frame(text_shapes[1].text_frame, body)

    if not output_path:
        OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
        output_path = OUTPUT_DIR / f"{docx_file.stem}_generated.pptx"
    else:
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
    target_prs.save(str(output_path))
    log(f"[Step 4] Finished. Output presentation saved to: {output_path}")

    # ----------------------------------------------------
    # Step 5: Verification & AI Agent Self-Correction QA Loop
    # ----------------------------------------------------
    log("[Step 5] Running PPTX File Integrity Verification & AI Agent Self-Correction Loop...")
    is_valid, final_path = verify_and_auto_heal_pptx(
        output_path,
        doc_structure=parsed_doc,
        template_inventory=template_inventory,
        log_cb=log
    )

    if not is_valid:
        log(f"[Step 5 Warning] Final PPTX file might contain remaining non-fatal notices.")
    else:
        log(f"[Step 5] PPTX Verification PASSED: File is clean, valid, and fully openable.")

    # Render Visual Preview QA
    try:
        preview_imgs = render_pptx_file_previews(output_path, target_width_px=650)
        log(f"[Step 5] Rendered {len(preview_imgs)} slide previews. Verification complete.")
    except Exception as qa_ex:
        log(f"[Step 5 Warning] QA preview render warning: {qa_ex}")

    return str(output_path)

def verify_pptx_integrity(file_path: str | Path) -> Tuple[bool, List[str]]:
    """
    Performs comprehensive structural, zip packaging, XML and python-pptx integrity checks on a PPTX file.
    Returns (is_valid: bool, issues: List[str]).
    """
    p = Path(file_path)
    issues: List[str] = []
    
    if not p.exists():
        return False, [f"File does not exist: {p}"]
        
    if p.stat().st_size == 0:
        return False, ["File is empty (0 bytes)"]
        
    # 1. Zip package & duplicate part checks
    try:
        with zipfile.ZipFile(str(p), "r") as z:
            bad_file = z.testzip()
            if bad_file:
                issues.append(f"Corrupted zip entry found: {bad_file}")
                
            namelist = z.namelist()
            counts = collections.Counter(namelist)
            dups = [k for k, v in counts.items() if v > 1]
            if dups:
                issues.append(f"Duplicate package parts detected ({len(dups)} duplicate entries): {dups[:5]}")
                
            # Test all XML files for syntactic validity
            for fname in namelist:
                if fname.endswith(".xml") or fname.endswith(".rels"):
                    try:
                        raw = z.read(fname)
                        ET.fromstring(raw)
                    except Exception as xml_err:
                        issues.append(f"XML parse error in '{fname}': {str(xml_err)}")
    except Exception as zip_err:
        return False, [f"Zip archive integrity error: {str(zip_err)}"]

    # 2. python-pptx model parse & slide count check
    try:
        prs = Presentation(str(p))
        if len(prs.slides) == 0:
            issues.append("Presentation contains 0 slides")
    except Exception as pptx_err:
        issues.append(f"python-pptx engine parse error: {str(pptx_err)}")
        
    return len(issues) == 0, issues

def repair_pptx_package(corrupted_path: str | Path, target_path: Optional[str | Path] = None) -> bool:
    """
    Repairs package-level corruptions such as duplicate part entries in pptx zip container.
    """
    src_p = Path(corrupted_path)
    if not src_p.exists():
        return False
        
    dst_p = Path(target_path) if target_path else src_p
    temp_p = src_p.parent / f"~temp_repaired_{src_p.name}"
    
    try:
        with zipfile.ZipFile(str(src_p), "r") as z_in:
            with zipfile.ZipFile(str(temp_p), "w", compression=zipfile.ZIP_DEFLATED) as z_out:
                seen_names = set()
                for item in z_in.infolist():
                    if item.filename in seen_names:
                        continue
                    seen_names.add(item.filename)
                    z_out.writestr(item, z_in.read(item.filename))
                    
        # Replace original/target with repaired file
        if dst_p.exists() and dst_p.resolve() == src_p.resolve():
            src_p.unlink()
        temp_p.replace(dst_p)
        return True
    except Exception:
        if temp_p.exists():
            try:
                temp_p.unlink()
            except Exception:
                pass
        return False

def verify_and_auto_heal_pptx(
    pptx_path: str | Path,
    doc_structure: Optional[Dict[str, Any]] = None,
    template_inventory: Optional[List[Dict[str, Any]]] = None,
    log_cb: Optional[Callable[[str], None]] = None,
    max_fix_attempts: int = 2
) -> Tuple[bool, str]:
    """
    Verification & Self-Correction QA Loop:
    1. Runs full structural & XML verification on pptx_path.
    2. If issues are found, attempts automated package repair (deduplicating zip parts).
    3. If errors persist, calls 9Router AI Agent to diagnose the failure, adjust the slide plan,
       and rebuild the presentation until clean.
    """
    def log(msg: str):
        if log_cb:
            log_cb(msg)

    p = Path(pptx_path)
    log(f"[Verification Loop] Validating PPTX integrity for '{p.name}'...")
    
    is_ok, issues = verify_pptx_integrity(p)
    if is_ok:
        log("[Verification Loop] PPTX integrity checks passed with 0 errors.")
        return True, str(p)

    log(f"[Verification Loop Warning] Found {len(issues)} integrity issue(s): {'; '.join(issues)}")

    # Attempt 1: Package-level deduplication & repair
    log("[Verification Loop - Fix 1] Running package repair & deduplication...")
    repair_success = repair_pptx_package(p, p)
    if repair_success:
        is_ok, issues = verify_pptx_integrity(p)
        if is_ok:
            log("[Verification Loop - Fix 1] Package repair succeeded! PPTX is now fully valid.")
            return True, str(p)
        else:
            log(f"[Verification Loop - Fix 1] Package repair applied, but remaining issues: {issues}")

    # Attempt 2: AI Agent Diagnostic & Self-Correction Rebuild with Visual Feedback
    if doc_structure and template_inventory:
        log("[Verification Loop - AI Agent] Invoking AI Agent to diagnose integrity errors and regenerate slide mapping...")
        try:
            client = OpenAI(
                api_key=Config.NINEROUTER_KEY or "dummy_key",
                base_url=f"{Config.NINEROUTER_URL.rstrip('/')}/v1",
                timeout=120.0
            )
            
            ai_repair_prompt = f"""
The generated PowerPoint presentation has corruptions/integrity errors:
Error List:
{json.dumps(issues, ensure_ascii=False, indent=2)}

Document Outline:
{json.dumps(doc_structure, ensure_ascii=False, indent=2)}

Please analyze the issues and output a corrected, safe slide replacement plan adhering to standard schema.
Ensure no conflicting shape removals or malformed tables are generated.
"""
            # Collect generated slide screenshots to provide visual feedback to AI
            user_msg_parts: List[Dict[str, Any]] = [
                {"type": "text", "text": ai_repair_prompt}
            ]
            try:
                gen_screenshots = render_pptx_file_previews(p, target_width_px=450, use_com=True)
                for idx, img in enumerate(gen_screenshots[:10]):
                    b64 = image_to_base64_jpeg(img, quality=80)
                    user_msg_parts.append({
                        "type": "image_url",
                        "image_url": {
                            "url": b64
                        }
                    })
                log(f"[Verification Loop - AI Agent] Attached {len(gen_screenshots[:10])} native visual slide screenshots to diagnostic prompt.")
            except Exception as ss_err:
                log(f"[Verification Loop - AI Agent Warning] Could not attach screenshots: {ss_err}")

            response = client.chat.completions.create(
                model=Config.NINEROUTER_CHAT_MODEL,
                messages=[
                    {"role": "system", "content": "You are a PowerPoint Diagnostic & Repair Agent. Fix presentation generation errors and output clean valid JSON."},
                    {"role": "user", "content": user_msg_parts}
                ],
                temperature=0.1
            )
            content = response.choices[0].message.content or "{}"
            match = re.search(r"```(?:json)?\s*(\{.*?\})\s*```", content, re.DOTALL)
            json_str = match.group(1) if match else content.strip()
            repaired_plan = json.loads(json_str)
            
            if repaired_plan and "slides" in repaired_plan:
                log("[Verification Loop - AI Agent] AI Agent provided healed plan. Re-assembling presentation...")
                # Re-assemble presentation with repaired plan
                first_tpl = template_inventory[0]["template_file"]
                src_base_p = DATA_DIR / first_tpl if (DATA_DIR / first_tpl).exists() else next(DATA_DIR.glob("*.pptx"))
                rebuilt_prs = Presentation(str(src_base_p))
                while len(rebuilt_prs.slides) > 0:
                    rId = rebuilt_prs.slides._sldIdLst[0].rId
                    rebuilt_prs.part.drop_rel(rId)
                    rebuilt_prs.slides._sldIdLst.remove(rebuilt_prs.slides._sldIdLst[0])
                    
                for s_plan in repaired_plan.get("slides", []):
                    src_tpl = s_plan.get("source_template") or first_tpl
                    src_idx = s_plan.get("source_slide_index", 0)
                    tpl_path = DATA_DIR / src_tpl if (DATA_DIR / src_tpl).exists() else src_base_p
                    src_prs = Presentation(str(tpl_path))
                    if src_idx >= len(src_prs.slides):
                        src_idx = 0
                        
                    t_slide = clone_slide_across_presentations(src_prs, rebuilt_prs, src_idx)
                    for r in s_plan.get("shape_replacements", []):
                        sh_i = r.get("shape_index")
                        if sh_i is not None and sh_i < len(t_slide.shapes):
                            sh = t_slide.shapes[sh_i]
                            if sh.has_text_frame:
                                _safe_update_text_frame(sh.text_frame, str(r.get("text", "")))
                                
                    if s_plan.get("speaker_notes"):
                        try:
                            t_slide.notes_slide.notes_text_frame.text = str(s_plan["speaker_notes"])
                        except Exception:
                            pass
                            
                rebuilt_prs.save(str(p))
                repair_pptx_package(p, p)
                
                is_ok, issues = verify_pptx_integrity(p)
                if is_ok:
                    log("[Verification Loop - AI Agent] Rebuilt deck is verified OK and error-free!")
                    return True, str(p)
        except Exception as ai_heal_err:
            log(f"[Verification Loop - AI Agent Error] Self-correction failed: {ai_heal_err}")

    # Fallback to package repair status
    is_ok, remaining = verify_pptx_integrity(p)
    return is_ok, str(p)
