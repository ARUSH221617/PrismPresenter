import json
import uuid
from pathlib import Path
from typing import Any, Dict, List, Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

from pptx_jahat.config import DATA_DIR, COMPONENTS_DIR, SHAPES_DIR, IMAGES_DIR
from pptx_jahat.tools.preview import render_pptx_slide_to_image, image_to_base64_jpeg

def _rgb_to_hex(color: Optional[RGBColor]) -> Optional[str]:
    if not color:
        return None
    try:
        return f"#{color[0]:02x}{color[1]:02x}{color[2]:02x}"
    except Exception:
        return None

def _get_shape_fill(shape: Any) -> Dict[str, Any]:
    fill_info = {"type": "none", "color": None}
    try:
        if shape.fill.type is not None:
            # Check solid fill
            if shape.fill.type == 1: # MSO_FILL.SOLID
                try:
                    fill_info["type"] = "solid"
                    fill_info["color"] = _rgb_to_hex(shape.fill.fore_color.rgb)
                except Exception:
                    fill_info["color"] = str(shape.fill.fore_color.theme_color) if hasattr(shape.fill.fore_color, 'theme_color') else None
    except Exception:
        pass
    return fill_info

def _get_shape_line(shape: Any) -> Dict[str, Any]:
    line_info = {"color": None, "width_pt": None}
    try:
        if shape.line and shape.line.fill.type is not None:
            try:
                line_info["color"] = _rgb_to_hex(shape.line.color.rgb)
            except Exception:
                pass
            if shape.line.width:
                line_info["width_pt"] = shape.line.width.pt
    except Exception:
        pass
    return line_info

def _get_shape_font(shape: Any) -> Dict[str, Any]:
    font_info = {"name": None, "size_pt": None, "bold": None, "italic": None, "color": None}
    try:
        if shape.has_text_frame and shape.text_frame.text:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.font:
                        if r.font.name: font_info["name"] = r.font.name
                        if r.font.size: font_info["size_pt"] = r.font.size.pt
                        if r.font.bold is not None: font_info["bold"] = r.font.bold
                        if r.font.italic is not None: font_info["italic"] = r.font.italic
                        try:
                            if r.font.color and r.font.color.rgb:
                                font_info["color"] = _rgb_to_hex(r.font.color.rgb)
                        except Exception:
                            pass
                        if font_info["name"] or font_info["size_pt"]:
                            return font_info
    except Exception:
        pass
    return font_info

def _determine_component_label_and_description(shape: Any, shape_type_name: str, text: str) -> tuple[str, str]:
    lower_text = text.lower().strip()
    words = lower_text.split()
    
    if shape.has_text_frame:
        font = _get_shape_font(shape)
        size = font.get("size_pt") or 14
        
        if size >= 24 or len(words) <= 6 and ("title" in lower_text or "header" in lower_text or size >= 20):
            return "Title / Header Box", "Prominent title or section header component"
        elif len(words) > 15 or "\n" in text:
            return "Content / Body Block", "Multi-line text block suitable for paragraphs and descriptions"
        elif len(words) <= 5 and any(char.isdigit() for char in text):
            return "Metric / Stat Callout", "Numeric statistic or key data highlight point"
        elif len(words) <= 4:
            return "Badge / Subtitle Label", "Short label, category badge, or subtitle"
        else:
            return "Text Box", "General purpose text container"
    
    if shape_type_name == "PICTURE":
        return "Image / Graphic Placeholder", "Visual container for photographs, illustrations or diagrams"
    elif shape_type_name == "TABLE":
        return "Data Table", "Structured grid for rows and columns of data"
    elif shape_type_name == "AUTO_SHAPE":
        return "Card / Container Shape", "Background card, container box, or geometric decoration"
    elif shape_type_name == "GROUP":
        return "Composite Group Card", "Grouped layout containing multiple combined visual elements"
    
    return f"{shape_type_name} Component", f"Presentation element of type {shape_type_name}"

def classify_slide_archetype(slide: Any, slide_idx: int, shapes_summary: List[Dict[str, Any]]) -> str:
    """
    Classifies a template slide into a semantic design archetype:
    - title_cover: Presentation master title/intro slide
    - chart_visual: Slide containing native charts/data visualizations
    - smartart_diagram: Flowcharts, process steps, hierarchies, cycles
    - table_matrix: Data tables & tabular comparisons
    - metrics_stats: Number callouts, KPIs, large statistics
    - multi_column: 2 or 3 column comparison / card layout
    - process_timeline: Sequential steps, numbered points, or flow
    - content_bullets: Standard topic overview with headings and bullet list
    - conclusion_quote: Summary, takeaway or final slide
    """
    if slide_idx == 0:
        return "title_cover"

    # Check for native charts & SmartArt
    has_chart = any(s.get("shape_type") == "CHART" or "chart" in s.get("shape_name", "").lower() for s in shapes_summary)
    if has_chart:
        return "chart_visual"

    has_smartart = any("smartart" in s.get("shape_name", "").lower() or "diagram" in s.get("shape_name", "").lower() for s in shapes_summary)
    if has_smartart:
        return "smartart_diagram"

    has_table = any(s.get("is_table") or s.get("shape_type") == "TABLE" for s in shapes_summary)
    if has_table:
        return "table_matrix"

    # Count metrics or numeric callouts
    numeric_callouts = 0
    body_blocks = 0
    cards_containers = 0
    for s in shapes_summary:
        txt = s.get("original_text", "")
        font_sz = s.get("font", {}).get("size_pt") or 14
        if any(c.isdigit() for c in txt) and font_sz >= 22 and len(txt.split()) <= 4:
            numeric_callouts += 1
        elif len(txt.split()) > 10:
            body_blocks += 1
        if "card" in s.get("shape_name", "").lower() or s.get("shape_type") in ["AUTO_SHAPE", "GROUP"]:
            cards_containers += 1

    if numeric_callouts >= 2:
        return "metrics_stats"
    if cards_containers >= 2 and body_blocks >= 2:
        return "multi_column"
    if any(kw in s.get("original_text", "").lower() for s in shapes_summary for kw in ["step", "phase", "مرحله", "راهبرد"]):
        return "process_timeline"
    if any(kw in s.get("original_text", "").lower() for s in shapes_summary for kw in ["thank", "summary", "پایان", "نتیجه", "conclusion"]):
        return "conclusion_quote"

    return "content_bullets"

def extract_pptx_file(pptx_path: Path) -> Dict[str, Any]:
    prs = Presentation(str(pptx_path))
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    extracted_slides = []
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_info = {
            "index": slide_idx,
            "layout_name": slide.slide_layout.name if slide.slide_layout else "Custom",
            "slide_width": slide_width,
            "slide_height": slide_height,
            "components": []
        }
        
        for shape_idx, shape in enumerate(slide.shapes):
            comp_id = f"{pptx_path.stem}_s{slide_idx}_c{shape_idx}_{uuid.uuid4().hex[:6]}"
            shape_type_name = str(shape.shape_type).replace("MSO_SHAPE_TYPE.", "")
            text_content = shape.text_frame.text if shape.has_text_frame else ""
            
            label, desc = _determine_component_label_and_description(shape, shape_type_name, text_content)
            
            comp_data: Dict[str, Any] = {
                "id": comp_id,
                "source_file": pptx_path.name,
                "slide_index": slide_idx,
                "shape_index": shape_idx,
                "type": shape_type_name,
                "label": label,
                "description": desc,
                "sample_text": text_content.strip()[:200] if text_content else "",
                "position": {
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height
                },
                "font": _get_shape_font(shape),
                "fill": _get_shape_fill(shape),
                "line": _get_shape_line(shape)
            }
            
            # Save images if shape is picture
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    img_filename = f"{comp_id}.{img.ext}"
                    img_path = IMAGES_DIR / img_filename
                    with open(img_path, "wb") as f:
                        f.write(img.blob)
                    comp_data["image_path"] = str(img_path)
                except Exception:
                    pass
            
            # Save table metadata
            if shape.has_table:
                comp_data["table_info"] = {
                    "rows": len(shape.table.rows),
                    "cols": len(shape.table.columns)
                }
                
            slide_info["components"].append(comp_data)
        
        extracted_slides.append(slide_info)
        
    return {
        "source_file": pptx_path.name,
        "slide_width": slide_width,
        "slide_height": slide_height,
        "total_slides": len(prs.slides),
        "slides": extracted_slides
    }

def extract_all_templates(data_dir: Optional[Path] = None) -> Dict[str, Any]:
    target_dir = data_dir or DATA_DIR
    components_file = COMPONENTS_DIR / "components.json"
    
    pptx_files = list(target_dir.glob("*.pptx"))
    catalog = {
        "templates": [],
        "all_components": [],
        "component_counts_by_label": {}
    }
    
    for pptx_file in pptx_files:
        try:
            data = extract_pptx_file(pptx_file)
            catalog["templates"].append(data)
            for s in data["slides"]:
                for c in s["components"]:
                    catalog["all_components"].append(c)
                    lbl = c["label"]
                    catalog["component_counts_by_label"][lbl] = catalog["component_counts_by_label"].get(lbl, 0) + 1
        except Exception as e:
            print(f"Error extracting {pptx_file}: {e}")
            
    with open(components_file, "w", encoding="utf-8") as f:
        json.dump(catalog, f, indent=2)
        
    return catalog

def get_components_catalog() -> Dict[str, Any]:
    components_file = COMPONENTS_DIR / "components.json"
    if components_file.exists():
        try:
            with open(components_file, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    return extract_all_templates()

def inspect_template_slides(pptx_path: Path | str, include_screenshots: bool = False, screenshot_width: int = 450) -> List[Dict[str, Any]]:
    """
    Detailed inspection of slide sequence and all text/table shapes for in-place cloning & editing.
    Optionally attaches base64 rendered slide screenshots.
    """
    path = Path(pptx_path)
    if not path.exists():
        return []
    
    prs = Presentation(str(path))
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    slides_summary = []
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_entry = {
            "template_file": path.name,
            "slide_index": slide_idx,
            "layout_name": slide.slide_layout.name if slide.slide_layout else f"Slide {slide_idx+1}",
            "slide_width": slide_w,
            "slide_height": slide_h,
            "text_slots": []
        }
        
        if include_screenshots:
            try:
                img = render_pptx_slide_to_image(slide, slide_w, slide_h, target_width_px=screenshot_width)
                slide_entry["screenshot_base64"] = image_to_base64_jpeg(img, quality=80)
            except Exception:
                slide_entry["screenshot_base64"] = None
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_type_name = str(shape.shape_type).replace("MSO_SHAPE_TYPE.", "")
            if shape.has_text_frame and shape.text_frame.text.strip():
                sample_text = shape.text_frame.text.strip()
                font_data = _get_shape_font(shape)
                slide_entry["text_slots"].append({
                    "shape_index": shape_idx,
                    "shape_name": shape.name,
                    "shape_type": shape_type_name,
                    "original_text": sample_text,
                    "font": font_data,
                    "is_title": (font_data.get("size_pt") or 14) >= 22 or "title" in shape.name.lower()
                })
            elif shape.has_table:
                table_cells = []
                for r_idx, row in enumerate(shape.table.rows):
                    row_texts = [cell.text.strip() for cell in row.cells]
                    table_cells.append(row_texts)
                slide_entry["text_slots"].append({
                    "shape_index": shape_idx,
                    "shape_name": shape.name,
                    "shape_type": "TABLE",
                    "is_table": True,
                    "table_rows": len(shape.table.rows),
                    "table_cols": len(shape.table.columns),
                    "original_table_data": table_cells
                })
            else:
                # Other non-text shapes (cards, icons, pictures) that AI may decide to remove or keep
                slide_entry["text_slots"].append({
                    "shape_index": shape_idx,
                    "shape_name": shape.name,
                    "shape_type": shape_type_name,
                    "is_decorative": True
                })
                
        # Classify semantic archetype
        slide_entry["archetype"] = classify_slide_archetype(slide, slide_idx, slide_entry["text_slots"])
        slides_summary.append(slide_entry)
        
    return slides_summary

def inspect_all_templates(data_dir: Optional[Path] = None, include_screenshots: bool = True, screenshot_width: int = 450) -> List[Dict[str, Any]]:
    """
    Scans and inspects all available PPTX templates in data folder.
    Returns a unified list of slide descriptors with shape slots and screenshot previews.
    """
    target_dir = data_dir or DATA_DIR
    pptx_files = sorted(list(target_dir.glob("*.pptx")))
    
    # Filter out generated presentations
    template_files = [f for f in pptx_files if not f.name.endswith("_generated.pptx")]
    
    all_slides = []
    for pptx_file in template_files:
        try:
            slides = inspect_template_slides(pptx_file, include_screenshots=include_screenshots, screenshot_width=screenshot_width)
            all_slides.extend(slides)
        except Exception as e:
            print(f"Warning: Failed to inspect template {pptx_file.name}: {e}")
            
    return all_slides
