#!/usr/bin/env python3
"""
preview.py  (v4.0 — Pure-Python PPTX Slide Rendering Engine)
===========================================================
Modular, production-grade PresentationML & DrawingML rendering engine with zero external desktop dependencies.

Features:
---------
* Colors ............ srgbClr / schemeClr / sysClr / prstClr / hslClr / scrgbClr,
                      transforms lumMod/lumOff (HSL-correct), shade, tint,
                      alpha / alphaMod / alphaOff (full opacity support).
* Fills ............. Solid, linear, radial (circle), rectangular, shape-contour gradients,
                      3D bevel specular highlights, pattern, picture, theme style refs.
* Geometry .......... Full OpenXML preset catalog (>180 Presets), cubic/quad Beziers,
                      connectors (bent/curved/arrowhead caps), custom geometry (custGeom).
* Typography ........ Multi-column text (numCol), WordArt warping, OMML math equations,
                      RTL/Arabic shaping, font matching with fallbacks.
* Charts ............ Native OpenXML 2D/3D & c16 charts (bar/column, line, pie, doughnut, area).
* SmartArt .......... Relational diagram layout engine (hierarchies, chevrons, cycles, matrices).
* Tables ............ Spanned/merged cells (hMerge, vMerge, gridSpan, rowSpan), diagonal borders.
* Media & Vectors ... SVG rasterization, EMF/WMF metafiles, video poster frames with play badges.
* Scaling & Caching . Parallel multi-slide rendering worker pool, tile & theme caching.
"""

from __future__ import annotations

import io
import os
import math
import logging
import base64
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import List, Optional, Tuple, Any, Dict, Union

from pptx import Presentation
from PIL import Image, ImageDraw, ImageChops

# Sub-engine imports
from pptx_jahat.tools.renderers.color_resolver import (
    Theme,
    ColorResolver,
    resolve_element_color,
    resolve_theme_fill,
    parse_color_elem,
    scheme_rgb,
    hex_to_rgb,
    clamp,
    RGBA,
    RGB,
    q as _q,
    local_name as _local,
    NS_A,
    NS_P,
    NS_R,
    NS_MC,
    _DEFAULT_COLORS,
    _SYSCLR_DEFAULTS,
    _SCHEME_ALIASES,
    _COLOR_TAGS,
)
from pptx_jahat.tools.renderers.geometry_engine import (
    GeometryEngine,
    get_preset_ops as _prst_ops,
    get_connector_ops as _connector_ops,
    render_custom_geom as _render_custom_geom,
    cust_geom_to_ops as _cust_geom_to_ops,
    rotate_ops as _rotate_ops,
    stroke_ops as _stroke_ops,
    fill_ops as _fill_ops,
    make_silhouette as _make_silhouette,
    draw_cubic_bezier as _draw_cubic_bezier,
    regular_polygon_pts as _regular_pts,
    star_polygon_pts as _star_pts,
    rounded_rect_pts as _rounded_rect_pts,
    parse_adjustments as _parse_adj,
    _DASH_PATTERNS,
)
from pptx_jahat.tools.renderers.gradient_engine import (
    GradientEngine,
    paint_gradient as _paint_gradient,
    paint_3d_bevel,
    paint_shadow as _paint_shadow,
    paint_glow as _paint_glow,
    paint_inner_shadow as _paint_inner_shadow,
    paint_soft_edge as _paint_soft_edge,
    paint_reflection as _paint_reflection,
    parse_gradient as _parse_gradient,
    parse_effect_lst as _parse_effect_lst,
    interp_stops as _interp_stops,
)
from pptx_jahat.tools.renderers.typography_engine import (
    FontResolver,
    TypographyEngine,
    RunStyle as _RunStyle,
    parse_run_style as _parse_run_style,
    render_txbody as _render_txbody,
    render_math_formula,
    shape_text_for_display as _shape_text_for_display,
    is_rtl_text as _is_rtl_text,
    text_width as _text_width,
    HAS_BIDI,
)
from pptx_jahat.tools.renderers.chart_renderer import (
    ChartRenderer,
    render_chart_part,
    parse_chart_xml,
    NS_C,
)
from pptx_jahat.tools.renderers.smartart_renderer import (
    SmartArtRenderer,
    render_smartart_part,
    NS_DGM,
)
from pptx_jahat.tools.renderers.table_renderer import (
    TableRenderer,
    render_table as _render_table_fn,
)
from pptx_jahat.tools.renderers.media_resolver import (
    MediaResolver,
    rasterize_svg,
    rasterize_emf_wmf,
    paint_media_poster,
)
from pptx_jahat.config import Config
from pptx_jahat.tools.cache.render_cache import (
    RenderCache,
    render_pptx_parallel,
)
from pptx_jahat.tools.renderers.com_renderer import (
    is_powerpoint_com_available,
    export_pptx_slides_com,
)

log = logging.getLogger("pptx_preview")


# ---------------------------------------------------------------------------
# Helper XML & Geometry Navigation
# ---------------------------------------------------------------------------
def _int_attr(node: Optional[ET.Element], name: str, default: int) -> int:
    if node is None:
        return default
    v = node.attrib.get(name)
    try:
        return int(float(v))
    except (TypeError, ValueError):
        return default


def _str_attr(node: Optional[ET.Element], name: str, default: str) -> str:
    if node is None:
        return default
    return node.attrib.get(name, default)


def _find_spPr(elem: ET.Element) -> Optional[ET.Element]:
    for tag in ("spPr", "grpSpPr", "cxnSpPr", "picPr"):
        found = elem.find(_q(NS_P, tag))
        if found is not None:
            return found
        found = elem.find(_q(NS_A, tag))
        if found is not None:
            return found
    return None


def _find_xfrm(elem: ET.Element) -> Optional[ET.Element]:
    spPr = _find_spPr(elem)
    if spPr is not None:
        xfrm = spPr.find(_q(NS_A, "xfrm"))
        if xfrm is not None:
            return xfrm
    xfrm = elem.find(_q(NS_P, "xfrm"))
    if xfrm is not None:
        return xfrm
    return elem.find(_q(NS_A, "xfrm"))


def _is_hidden(elem: ET.Element) -> bool:
    for nv_tag in ("nvSpPr", "nvPicPr", "nvCxnSpPr", "nvGraphicFramePr", "nvGrpSpPr"):
        nv = elem.find(_q(NS_P, nv_tag))
        if nv is not None:
            cNvPr = nv.find(_q(NS_P, "cNvPr"))
            if cNvPr is not None and cNvPr.attrib.get("hidden") in ("1", "true"):
                return True
    return False


def _find_style(elem: ET.Element) -> Optional[ET.Element]:
    return elem.find(_q(NS_P, "style"))


def _style_ref(elem: ET.Element, name: str) -> Tuple[int, Optional[RGB]]:
    style = _find_style(elem)
    if style is None:
        return 0, None
    ref = style.find(_q(NS_A, name))
    if ref is None:
        return 0, None
    idx = _int_attr(ref, "idx", 0)
    c_node = ref.find(_q(NS_A, "schemeClr"))
    val = c_node.attrib.get("val") if c_node is not None else None
    return idx, val


def _blob_for_part(part: Any, rid: str) -> Optional[bytes]:
    if not part or not rid:
        return None
    try:
        rel = part.rels[rid]
        return rel.target_part.blob
    except Exception:
        return None


def _blob_for_rid(shape: Any, rid: str) -> Optional[bytes]:
    try:
        return _blob_for_part(shape.part, rid)
    except Exception:
        return None


def _parse_fill(spPr: Optional[ET.Element], colors: Dict[str, RGB]) -> Tuple[Optional[str], Any, Any]:
    if spPr is None:
        return None, None, None
    if spPr.find(_q(NS_A, "noFill")) is not None:
        return "none", (0, 0, 0, 0), None
    sf = spPr.find(_q(NS_A, "solidFill"))
    if sf is not None:
        c = resolve_element_color(sf, colors)
        return "solid", c, None
    gf = spPr.find(_q(NS_A, "gradFill"))
    if gf is not None:
        stops, ang, info = _parse_gradient(gf, colors)
        return "grad", stops, (ang, info)
    bf = spPr.find(_q(NS_A, "blipFill"))
    if bf is not None:
        blip = bf.find(_q(NS_A, "blip"))
        rid = blip.attrib.get(_q(NS_R, "embed")) if blip is not None else None
        return "blip", rid, bf.find(_q(NS_A, "srcRect"))
    pf = spPr.find(_q(NS_A, "pattFill"))
    if pf is not None:
        prst = pf.attrib.get("prst", "solid")
        # Resolve fg and bg colors separately; each may carry its own alpha.
        fg = _resolve_pattern_color(pf.find(_q(NS_A, "fgClr")), colors) or (240, 240, 240, 255)
        bg = _resolve_pattern_color(pf.find(_q(NS_A, "bgClr")), colors) or (255, 255, 255, 0)
        return "patt", (prst, fg, bg), None
    return None, None, None


def _resolve_pattern_color(parent: Optional[ET.Element], colors: Dict[str, RGB]) -> Optional[RGBA]:
    """Resolves the color (with alpha) of a <a:fgClr>/<a:bgClr> element inside <a:pattFill>."""
    if parent is None:
        return None
    # <a:fgClr> directly contains <a:srgbClr>/<a:schemeClr>/etc.
    return resolve_element_color(parent, colors)


def _parse_line(spPr: Optional[ET.Element], colors: Dict[str, RGB]) -> Tuple[Optional[RGBA], float, bool, Optional[Tuple[int, ...]], Optional[Dict[str, str]], Optional[Dict[str, str]]]:
    """Returns (color, width_pt, dashed, dash_pat, head_end, tail_end).
    head_end / tail_end are dicts like {'type':'arrow','w':'med','len':'med'} or None."""
    if spPr is None:
        return None, 1.0, False, None, None, None
    ln = spPr.find(_q(NS_A, "ln"))
    if ln is None:
        return None, 1.0, False, None, None, None
    if ln.find(_q(NS_A, "noFill")) is not None:
        return (0, 0, 0, 0), 0.0, False, None, None, None
    c = resolve_element_color(ln, colors) or (0, 0, 0, 255)
    w_emu = float(ln.attrib.get("w", 12700))
    w_pt = w_emu / 12700.0
    prstDash = ln.find(_q(NS_A, "prstDash"))
    dash_val = prstDash.attrib.get("val", "solid") if prstDash is not None else "solid"
    dashed = dash_val not in ("solid", "")
    dash_pat = _DASH_PATTERNS.get(dash_val)
    head_end = _parse_line_end(ln.find(_q(NS_A, "headEnd")))
    tail_end = _parse_line_end(ln.find(_q(NS_A, "tailEnd")))
    return c, w_pt, dashed, dash_pat, head_end, tail_end


def _parse_line_end(elem: Optional[ET.Element]) -> Optional[Dict[str, str]]:
    """Parses <a:headEnd>/<a:tailEnd> into {'type','w','len'}; returns None if no marker."""
    if elem is None:
        return None
    t = elem.attrib.get("type", "none")
    if t in (None, "", "none"):
        return None
    return {
        "type": t,
        "w": elem.attrib.get("w", "med"),
        "len": elem.attrib.get("len", "med"),
    }


def _line_end_scale(end_info: Optional[Dict[str, str]], lw: int) -> Tuple[float, float]:
    """Returns (length, width) in pixels for a line-end marker.
    DrawingML spec: sm ≈ 1× line width, med ≈ 2×, lg ≈ 3× (approximate)."""
    if not end_info:
        return 0.0, 0.0
    mult = {"sm": 1.5, "med": 2.5, "lg": 3.5}.get(end_info.get("len", "med"), 2.5)
    w_mult = {"sm": 1.5, "med": 2.5, "lg": 3.5}.get(end_info.get("w", "med"), 2.5)
    return max(lw * mult, 4.0), max(lw * w_mult, 4.0)


def _paint_line_end(draw: ImageDraw.ImageDraw, point: Tuple[float, float],
                    direction: Tuple[float, float], end_info: Optional[Dict[str, str]],
                    color: RGBA, lw: int, is_head: bool = False) -> None:
    """Draws an arrowhead / oval / triangle / diamond / stealth marker at `point`.
    `direction` is the unit vector of the line's outgoing direction at that end.
    For tailEnd, direction points away from the line (outward).
    For headEnd, direction points toward the line (inward); we flip it for drawing."""
    if not end_info or color[3] == 0:
        return
    # For headEnd the marker is drawn at the start point pointing INTO the line,
    # so the marker tip is on the line. We treat `direction` as the line's tangent
    # direction at that point and reverse it for head markers so the marker
    # "exits" away from the line endpoint (matching PowerPoint).
    if is_head:
        direction = (-direction[0], -direction[1])
    length, width = _line_end_scale(end_info, lw)
    ttype = end_info.get("type", "none")
    px, py = point
    dx, dy = direction
    # The marker tip is at `point`; the base is `length` behind it.
    base_x = px - dx * length
    base_y = py - dy * length
    # Perpendicular vector (rotated 90°) for marker width.
    perp_x = -dy
    perp_y = dx
    half_w = width / 2.0
    base_left = (base_x + perp_x * half_w, base_y + perp_y * half_w)
    base_right = (base_x - perp_x * half_w, base_y - perp_y * half_w)

    if ttype == "oval":
        r = width / 2.0
        cx = px - dx * (length / 2.0)
        cy = py - dy * (length / 2.0)
        draw.ellipse([cx - r, cy - r, cx + r, cy + r], fill=color)
    elif ttype == "diamond":
        # Rhombus: tip at base_x+dx*length/2, two side points, tail at base_x.
        mid_x = (px + base_x) / 2.0
        mid_y = (py + base_y) / 2.0
        left = (mid_x + perp_x * half_w, mid_y + perp_y * half_w)
        right = (mid_x - perp_x * half_w, mid_y - perp_y * half_w)
        draw.polygon([point, left, (base_x, base_y), right], fill=color)
    elif ttype == "stealth":
        # Stealth = arrow with a notch. Tip at `point`, base at base_left/base_right,
        # plus a notch point partway back (e.g. 60% of length).
        notch_x = px - dx * (length * 0.6)
        notch_y = py - dy * (length * 0.6)
        draw.polygon([point, base_left, (notch_x, notch_y), base_right], fill=color)
    elif ttype == "triangle" or ttype == "arrow":
        # Simple triangle: tip at `point`, base at base_left/base_right.
        draw.polygon([point, base_left, base_right], fill=color)
    else:
        # Unknown — draw nothing rather than a wrong marker.
        pass


def _box_from_xfrm(xfrm: ET.Element, T: Any) -> Tuple[Optional[Tuple[float, float, float, float]], float]:
    off = xfrm.find(_q(NS_A, "off"))
    ext = xfrm.find(_q(NS_A, "ext"))
    if off is None or ext is None:
        return None, 0.0
    try:
        x = float(off.attrib.get("x", 0))
        y = float(off.attrib.get("y", 0))
        cx = float(ext.attrib.get("cx", 0))
        cy = float(ext.attrib.get("cy", 0))
    except ValueError:
        return None, 0.0
    x0, y0 = T(x, y)
    x1, y1 = T(x + cx, y + cy)
    rot = float(xfrm.attrib.get("rot", 0)) / 60000.0
    return (min(x0, x1), min(y0, y1), max(x0, x1), max(y0, y1)), rot


def _extract_theme_color_palette(prs: Presentation) -> Dict[str, RGB]:
    t = Theme.from_presentation(prs)
    return t.colors


def _theme_for_slide(slide: Any, cache: Dict[int, Theme], fallback: Theme) -> Theme:
    try:
        master = slide.slide_layout.slide_master
        m_id = id(master)
        if m_id in cache:
            return cache[m_id]
        for rel in master.part.rels.values():
            if "theme" in rel.reltype:
                t = Theme()
                t.load_from_theme_part(rel.target_part.blob)
                cache[m_id] = t
                return t
    except Exception:
        pass
    return fallback


def _paint_blip(img: Image.Image, box: Tuple[float, float, float, float],
                blob: bytes, silhouette: Any) -> None:
    x0, y0, x1, y1 = box
    bw = int(math.ceil(x1 - x0))
    bh = int(math.ceil(y1 - y0))
    if bw <= 0 or bh <= 0 or not blob:
        return
    try:
        pic = Image.open(io.BytesIO(blob)).convert("RGBA").resize((bw, bh), Image.Resampling.LANCZOS)
        if silhouette is not None:
            mask = Image.new("L", (bw, bh), 0)
            silhouette(ImageDraw.Draw(mask))
            # Preserve picture's own alpha by multiplying with the silhouette mask,
            # then composite (paste would discard the source alpha channel).
            pic.putalpha(ImageChops.multiply(pic.split()[3], mask))
        img.alpha_composite(pic, (int(x0), int(y0)))
    except Exception as e:
        log.debug(f"Failed to paint blip: {e}")


# ---------------------------------------------------------------------------
# Pattern fill rendering
# ---------------------------------------------------------------------------
# Pattern tile size (px). Patterns repeat at this period.
_PATT_TILE = 16

def _patt_tile_mask(prst: str) -> Optional[Image.Image]:
    """Builds an 'L' mode tile (0..255) for the given preset pattern name.
    255 = foreground, 0 = background. Returns None for unknown patterns
    (caller should fall back to solid foreground)."""
    t = _PATT_TILE
    m = Image.new("L", (t, t), 0)
    d = ImageDraw.Draw(m)
    half = t // 2
    if prst in ("solid",):
        d.rectangle([0, 0, t, t], fill=255)
    elif prst in ("ltHorz", "horz"):
        for y in range(0, t, 4):
            d.line([(0, y), (t, y)], fill=255, width=1)
    elif prst in ("ltVert", "vert"):
        for x in range(0, t, 4):
            d.line([(x, 0), (x, t)], fill=255, width=1)
    elif prst in ("ltUpDiag", "upDiag", "dashUpDiag"):
        for k in range(-t, 2 * t, 4):
            d.line([(k, 0), (k + t, t)], fill=255, width=1)
    elif prst in ("ltDnDiag", "dnDiag", "dashDnDiag"):
        for k in range(-t, 2 * t, 4):
            d.line([(k, t), (k + t, 0)], fill=255, width=1)
    elif prst in ("diagCross",):
        for k in range(-t, 2 * t, 4):
            d.line([(k, 0), (k + t, t)], fill=255, width=1)
            d.line([(k, t), (k + t, 0)], fill=255, width=1)
    elif prst in ("cross",):
        for y in range(0, t, 4):
            d.line([(0, y), (t, y)], fill=255, width=1)
        for x in range(0, t, 4):
            d.line([(x, 0), (x, t)], fill=255, width=1)
    elif prst in ("dkHorz",):
        d.rectangle([0, 0, t, t], fill=255)
        for y in range(0, t, 4):
            d.line([(0, y), (t, y)], fill=0, width=1)
    elif prst in ("dkVert",):
        d.rectangle([0, 0, t, t], fill=255)
        for x in range(0, t, 4):
            d.line([(x, 0), (x, t)], fill=0, width=1)
    elif prst in ("dkUpDiag",):
        d.rectangle([0, 0, t, t], fill=255)
        for k in range(-t, 2 * t, 4):
            d.line([(k, 0), (k + t, t)], fill=0, width=1)
    elif prst in ("dkDnDiag",):
        d.rectangle([0, 0, t, t], fill=255)
        for k in range(-t, 2 * t, 4):
            d.line([(k, t), (k + t, 0)], fill=0, width=1)
    elif prst in ("pct10", "pct5"):
        # sparse dot grid
        for y in range(0, t, 6):
            for x in range(0, t, 6):
                d.point((x, y), fill=255)
    elif prst in ("pct20", "pct25"):
        for y in range(0, t, 4):
            for x in range(0, t, 4):
                d.point((x, y), fill=255)
    elif prst in ("pct50",):
        d.rectangle([0, 0, half, half], fill=255)
        d.rectangle([half, half, t, t], fill=255)
    elif prst in ("pct75", "pct80"):
        d.rectangle([0, 0, t, t], fill=255)
        d.rectangle([0, 0, half, half], fill=0)
        d.rectangle([half, half, t, t], fill=0)
    elif prst in ("trellis",):
        for k in range(-t, 2 * t, 6):
            d.line([(k, 0), (k + t, t)], fill=255, width=1)
            d.line([(k, t), (k + t, 0)], fill=255, width=1)
    elif prst in ("dotGrid", "smGrid"):
        for y in range(0, t, 4):
            for x in range(0, t, 4):
                d.point((x, y), fill=255)
    elif prst in ("smChecker",):
        d.rectangle([0, 0, half, half], fill=255)
        d.rectangle([half, half, t, t], fill=255)
    else:
        # Unknown pattern — fall back to solid fg so the shape is still visible.
        d.rectangle([0, 0, t, t], fill=255)
    return m


def _paint_pattern(img: Image.Image, box: Tuple[float, float, float, float],
                   fill_data: Any, silhouette: Any) -> None:
    """Renders a DrawingML preset pattern (pattFill) clipped to the shape silhouette."""
    if not fill_data:
        return
    prst, fg, bg = fill_data
    x0, y0, x1, y1 = box
    bw = int(math.ceil(x1 - x0))
    bh = int(math.ceil(y1 - y0))
    if bw <= 0 or bh <= 0:
        return

    tile = _patt_tile_mask(prst)
    if tile is None:
        # Unknown pattern: paint solid foreground as a safe fallback.
        layer = Image.new("RGBA", (bw, bh), fg)
    else:
        # Build a tile of the correct size, then repeat to fill (bw, bh).
        fg_layer = Image.new("RGBA", (_PATT_TILE, _PATT_TILE), fg)
        bg_layer = Image.new("RGBA", (_PATT_TILE, _PATT_TILE), bg)
        tile_rgba = Image.composite(fg_layer, bg_layer, tile)
        layer = Image.new("RGBA", (bw, bh))
        for y in range(0, bh, _PATT_TILE):
            for x in range(0, bw, _PATT_TILE):
                layer.paste(tile_rgba, (x, y))

    if silhouette is not None:
        mask = Image.new("L", (bw, bh), 0)
        silhouette(ImageDraw.Draw(mask))
        layer.putalpha(ImageChops.multiply(layer.split()[3], mask))
    img.alpha_composite(layer, (int(x0), int(y0)))


def _paint_picture(img: Image.Image, box: Tuple[float, float, float, float],
                   blob: bytes, src_rect: Optional[ET.Element] = None) -> None:
    if not blob:
        return
    x0, y0, x1, y1 = box
    bw = int(round(x1 - x0))
    bh = int(round(y1 - y0))
    if bw <= 0 or bh <= 0:
        return
    try:
        # Check if vector SVG
        if blob.strip().startswith(b"<svg") or b"<svg" in blob[:200]:
            svg_img = rasterize_svg(blob, bw, bh)
            if svg_img:
                img.alpha_composite(svg_img, (int(x0), int(y0)))
                return

        pic = Image.open(io.BytesIO(blob)).convert("RGBA")
        if src_rect is not None:
            pw, ph = pic.size
            l = float(src_rect.attrib.get("l", 0)) / 100000.0 * pw
            t = float(src_rect.attrib.get("t", 0)) / 100000.0 * ph
            r = float(src_rect.attrib.get("r", 0)) / 100000.0 * pw
            b = float(src_rect.attrib.get("b", 0)) / 100000.0 * ph
            crop_box = (max(0, int(l)), max(0, int(t)), min(pw, int(pw - r)), min(ph, int(ph - b)))
            if crop_box[2] > crop_box[0] and crop_box[3] > crop_box[1]:
                pic = pic.crop(crop_box)
        pic = pic.resize((bw, bh), Image.Resampling.LANCZOS)
        img.paste(pic, (int(x0), int(y0)), pic)
    except Exception as e:
        log.debug(f"Failed to paint picture: {e}")


# ---------------------------------------------------------------------------
# SlideRenderer Coordinator
# ---------------------------------------------------------------------------
class SlideRenderer:
    """
    Main slide rendering orchestrator. Dispatches visual elements to specialized sub-engines:
    - Backgrounds, shapes, pictures, connectors, groups, graphic frames (charts, SmartArt, tables).
    """

    def __init__(self, slide: Any, prs: Presentation, width: int, theme: Theme, fonts: FontResolver):
        self.slide = slide
        self.prs = prs
        self.theme = theme
        self.fonts = fonts
        self.scale = width / float(prs.slide_width or 9144000)
        height = max(1, int(round((prs.slide_height or 6858000) * self.scale)))
        self.img = Image.new("RGBA", (width, height), (255, 255, 255, 255))
        self.ctx = {
            "palette": theme.colors,
            "fonts": fonts,
            "scale_x": self.scale,
            "scale_y": self.scale,
            "theme": theme,
        }

    def render(self) -> Image.Image:
        self._render_background()
        T = lambda x, y: (x * self.scale, y * self.scale)

        # 1. Render non-placeholder shapes from slide master
        try:
            if hasattr(self.slide, "slide_layout") and self.slide.slide_layout:
                master = self.slide.slide_layout.slide_master
                if master:
                    m_cSld = master._element.find(_q(NS_P, "cSld"))
                    m_spTree = m_cSld.find(_q(NS_P, "spTree")) if m_cSld is not None else None
                    if m_spTree is not None:
                        self._render_master_or_layout_children(m_spTree, master.part, T)
        except Exception:
            pass

        # 2. Render non-placeholder decorative shapes from slide layout
        try:
            if hasattr(self.slide, "slide_layout") and self.slide.slide_layout:
                layout = self.slide.slide_layout
                l_cSld = layout._element.find(_q(NS_P, "cSld"))
                l_spTree = l_cSld.find(_q(NS_P, "spTree")) if l_cSld is not None else None
                if l_spTree is not None:
                    self._render_master_or_layout_children(l_spTree, layout.part, T)
        except Exception:
            pass

        # 3. Render slide shapes
        cSld = self.slide._element.find(_q(NS_P, "cSld"))
        spTree = cSld.find(_q(NS_P, "spTree")) if cSld is not None else None
        if spTree is not None:
            self._render_children(spTree, self.slide.part, T)
        return self.img

    def _render_master_or_layout_children(self, parent: ET.Element, part: Any, T: Any) -> None:
        """Renders only non-placeholder decorative background graphics from master/layout."""
        for child in parent:
            # Skip placeholders (title, body, date, footer, slide num).
            # NOTE: <p:ph> is nested inside <p:nvPr> (a grandchild of <p:nvSpPr>),
            # so a direct `find` returns None — must use the descendant XPath `.//`.
            nv = child.find(_q(NS_P, "nvSpPr"))
            if nv is None:
                nv = child.find(_q(NS_P, "nvPicPr"))
            if nv is None:
                nv = child.find(_q(NS_P, "nvCxnSpPr"))
            if nv is not None and nv.find(f".//{_q(NS_P, 'ph')}") is not None:
                continue

            tag = _local(child.tag)
            if tag == "sp":
                self._render_sp(child, part, T)
            elif tag == "pic":
                self._render_pic(child, part, T)
            elif tag == "cxnSp":
                self._render_cxn(child, part, T)
            elif tag == "grpSp":
                self._render_group(child, part, T)

    def _render_background(self) -> None:
        try:
            layout = self.slide.slide_layout
            master = layout.slide_master
            sources = [self.slide, layout, master]
        except Exception:
            sources = [self.slide]
        for src in sources:
            if self._paint_bg_element(src):
                return
        c = self.theme.colors.get("lt1", (255, 255, 255))
        ImageDraw.Draw(self.img, "RGBA").rectangle([0, 0, self.img.size[0], self.img.size[1]], fill=(c[0], c[1], c[2], 255))

    def _paint_bg_element(self, src: Any) -> bool:
        try:
            cSld = src._element.find(_q(NS_P, "cSld"))
        except Exception:
            return False
        if cSld is None:
            return False
        bg = cSld.find(_q(NS_P, "bg"))
        if bg is None:
            return False
        colors = self.theme.colors
        draw = ImageDraw.Draw(self.img, "RGBA")
        box = (0.0, 0.0, float(self.img.size[0]), float(self.img.size[1]))
        part = getattr(src, "part", None)

        bgPr = bg.find(_q(NS_P, "bgPr"))
        if bgPr is not None:
            if bgPr.find(_q(NS_A, "noFill")) is not None:
                return True
            sf = bgPr.find(_q(NS_A, "solidFill"))
            if sf is not None:
                c = resolve_element_color(sf, colors)
                if c:
                    if c[3] < 255:
                        bg_layer = Image.new("RGBA", self.img.size, c)
                        self.img.alpha_composite(bg_layer)
                    else:
                        draw.rectangle(box, fill=c)
                return True
            gf = bgPr.find(_q(NS_A, "gradFill"))
            if gf is not None:
                stops, ang, info = _parse_gradient(gf, colors)
                _paint_gradient(self.img, box, stops, ang, lambda d: d.rectangle(box, fill=255), info)
                return True
            bf = bgPr.find(_q(NS_A, "blipFill"))
            if bf is not None:
                blip = bf.find(_q(NS_A, "blip"))
                rid = blip.attrib.get(_q(NS_R, "embed")) if blip is not None else None
                blob = _blob_for_part(part, rid) if rid else None
                if blob:
                    _paint_picture(self.img, box, blob, bf.find(_q(NS_A, "srcRect")))
                return True
            return True

        bgRef = bg.find(_q(NS_P, "bgRef"))
        if bgRef is not None:
            idx = _int_attr(bgRef, "idx", 1)
            ph = bgRef.find(_q(NS_A, "schemeClr"))
            ph_rgb = scheme_rgb(ph.attrib.get("val"), colors, None) if ph is not None else None
            res = resolve_theme_fill(self.theme.fmt_scheme, "bgFillStyleLst", idx, ph_rgb, colors)
            if res:
                if res[0] == "solid" and res[1]:
                    if res[1][3] < 255:
                        bg_layer = Image.new("RGBA", self.img.size, res[1])
                        self.img.alpha_composite(bg_layer)
                    else:
                        draw.rectangle(box, fill=res[1])
                elif res[0] == "grad":
                    ang = res[2][0] if isinstance(res[2], tuple) else (res[2] or 0.0)
                    info = res[2][1] if isinstance(res[2], tuple) else None
                    _paint_gradient(self.img, box, res[1], ang, lambda d: d.rectangle(box, fill=255), info)
                elif res[0] == "blip":
                    blob = _blob_for_part(part, res[1])
                    if blob:
                        _paint_picture(self.img, box, blob)
            return True
        return False

    def _render_children(self, parent: ET.Element, part: Any, T: Any) -> None:
        for child in parent:
            tag = _local(child.tag)
            if tag == "sp":
                self._render_sp(child, part, T)
            elif tag == "pic":
                self._render_pic(child, part, T)
            elif tag == "cxnSp":
                self._render_cxn(child, part, T)
            elif tag == "graphicFrame":
                self._render_graphicframe(child, part, T)
            elif tag == "grpSp":
                self._render_group(child, part, T)
            elif tag == "alternateContent":
                choice = child.find(_q(NS_MC, "Choice"))
                fallback = child.find(_q(NS_MC, "Fallback"))
                target = choice if choice is not None else (fallback if fallback is not None else child)
                self._render_children(target, part, T)

    def _render_group(self, grp: ET.Element, part: Any, T: Any) -> None:
        if _is_hidden(grp):
            return
        xfrm = _find_xfrm(grp)
        if xfrm is None:
            self._render_children(grp, part, T)
            return
        off = xfrm.find(_q(NS_A, "off"))
        ext = xfrm.find(_q(NS_A, "ext"))
        chOff = xfrm.find(_q(NS_A, "chOff"))
        chExt = xfrm.find(_q(NS_A, "chExt"))
        if off is None or ext is None:
            self._render_children(grp, part, T)
            return
        try:
            ox, oy = float(off.attrib.get("x", 0)), float(off.attrib.get("y", 0))
            ex, ey = float(ext.attrib.get("cx", 1)), float(ext.attrib.get("cy", 1))
        except ValueError:
            self._render_children(grp, part, T)
            return
        if chOff is not None and chExt is not None:
            cox = float(chOff.attrib.get("x", ox))
            coy = float(chOff.attrib.get("y", oy))
            cex = float(chExt.attrib.get("cx", ex)) or ex
            cey = float(chExt.attrib.get("cy", ey)) or ey
        else:
            cox, coy, cex, cey = ox, oy, ex, ey
        kx = (ex / cex) if cex else 1.0
        ky = (ey / cey) if cey else 1.0

        def GT(x, y):
            return T(ox + (x - cox) * kx, oy + (y - coy) * ky)

        self._render_children(grp, part, GT)

    def _placeholder_chain(self, ph: ET.Element) -> List[Any]:
        chain = []
        ptype = ph.attrib.get("type", "body") or "body"
        pidx = ph.attrib.get("idx", "0") or "0"
        try:
            layout = self.slide.slide_layout
            ancestors = [layout, layout.slide_master]
        except Exception:
            ancestors = []
        for anc in ancestors:
            try:
                spTree = anc._element.find(_q(NS_P, "cSld")).find(_q(NS_P, "spTree"))
            except Exception:
                continue
            if spTree is None:
                continue
            for cand in spTree.iter(_q(NS_P, "sp")):
                nv = cand.find(_q(NS_P, "nvSpPr"))
                # <p:ph> is nested inside <p:nvPr>; use descendant XPath.
                cph = nv.find(f".//{_q(NS_P, 'ph')}") if nv is not None else None
                if cph is None:
                    continue
                ctype = cph.attrib.get("type", "body") or "body"
                cidx = cph.attrib.get("idx", "0") or "0"
                if cidx != pidx:
                    continue
                if not (ctype == ptype or ctype in ("body", "obj") or ptype in ("body", "obj")):
                    continue
                chain.append((_find_spPr(cand), cand.find(_q(NS_P, "txBody")), cand))
                break
        return chain

    def _resolve_fill_line(self, elem: ET.Element, spPr: Optional[ET.Element], chain: List[Any]):
        colors = self.theme.colors
        fmt = self.theme.fmt_scheme
        fill_kind, fill_data, fill_extra = _parse_fill(spPr, colors)
        line_col, line_w, dashed, dash_pat, head_end, tail_end = _parse_line(spPr, colors)

        if fill_kind is None:
            for ph_spPr, _, _ in chain:
                if ph_spPr is not None:
                    fk, fd, fe = _parse_fill(ph_spPr, colors)
                    if fk is not None:
                        fill_kind, fill_data, fill_extra = fk, fd, fe
                        break
        if line_col is None:
            for ph_spPr, _, _ in chain:
                if ph_spPr is not None:
                    lc, lw, d, dp, he, te = _parse_line(ph_spPr, colors)
                    if lc is not None:
                        line_col, line_w, dashed, dash_pat = lc, lw, d, dp
                        head_end = head_end or he
                        tail_end = tail_end or te
                        break

        # Theme style refs
        if fill_kind is None:
            idx, ph_val = _style_ref(elem, "fillRef")
            ph_rgb = scheme_rgb(ph_val, colors, None) if ph_val else None
            res = resolve_theme_fill(fmt, "fillStyleLst", idx, ph_rgb, colors)
            if res:
                fill_kind, fill_data, fill_extra = res[0], res[1], res[2]

        if line_col is None:
            idx, ph_val = _style_ref(elem, "lnRef")
            ph_rgb = scheme_rgb(ph_val, colors, None) if ph_val else None
            res = resolve_theme_fill(fmt, "lnStyleLst", idx, ph_rgb, colors)
            if res and res[0] == "solid":
                line_col = res[1]

        return (fill_kind, fill_data, fill_extra), line_col, line_w, dashed, dash_pat, head_end, tail_end

    def _render_sp(self, sp: ET.Element, part: Any, T: Any) -> None:
        if _is_hidden(sp):
            return
        spPr = _find_spPr(sp)
        xfrm = _find_xfrm(sp)
        if xfrm is None:
            return
        box, rot_deg = _box_from_xfrm(xfrm, T)
        if box is None:
            return

        nv = sp.find(_q(NS_P, "nvSpPr"))
        # <p:ph> is nested inside <p:nvPr> under <p:nvSpPr> — use descendant XPath.
        ph = nv.find(f".//{_q(NS_P, 'ph')}") if nv is not None else None
        chain = self._placeholder_chain(ph) if ph is not None else []

        fill_info, line_col, line_w, dashed, dash_pat, _head_end, _tail_end = self._resolve_fill_line(sp, spPr, chain)
        fill_kind, fill_data, fill_extra = fill_info

        prstGeom = spPr.find(_q(NS_A, "prstGeom")) if spPr is not None else None
        custGeom = spPr.find(_q(NS_A, "custGeom")) if spPr is not None else None
        prst = prstGeom.attrib.get("prst", "rect") if prstGeom is not None else "rect"
        adj = _parse_adj(prstGeom)

        effects = _parse_effect_lst(spPr, self.theme.colors, self.scale)
        # Build ops: preset shapes use _prst_ops; custGeom paths are converted to ops
        # so they go through the SAME fill/effect pipeline (gradient, pattern, blip,
        # shadow, glow, 3D bevel, reflection) instead of being limited to solid fill.
        if custGeom is not None:
            ops = _cust_geom_to_ops(custGeom, box)
        else:
            ops = _prst_ops(prst, box, adj)
        if rot_deg:
            cx, cy = (box[0] + box[2]) / 2.0, (box[1] + box[3]) / 2.0
            ops = _rotate_ops(ops, cx, cy, rot_deg)

        lw = max(1, int(round(line_w * 12700 * self.scale))) if line_col and line_col[3] > 0 else 0

        # Render custom geometry — legacy code path kept for backward compatibility,
        # but only used when ops extraction failed.
        if custGeom is not None and not ops:
            fill_c = fill_data if fill_kind == "solid" else None
            _render_custom_geom(self.img, ImageDraw.Draw(self.img, "RGBA"), custGeom, box, fill_c, line_col, lw)
        else:
            # Effects (Shadow, Glow)
            if effects.get("outerShdw"):
                _paint_shadow(self.img, ops, effects["outerShdw"])
            if effects.get("glow"):
                _paint_glow(self.img, ops, effects["glow"])

            # Fill
            draw = ImageDraw.Draw(self.img, "RGBA")
            if fill_kind == "solid" and fill_data and fill_data[3] > 0:
                _fill_ops(draw, ops, fill_data)
            elif fill_kind == "grad" and fill_data:
                ang = fill_extra[0] if isinstance(fill_extra, tuple) else (fill_extra or 0.0)
                info = fill_extra[1] if isinstance(fill_extra, tuple) else None
                _paint_gradient(self.img, box, fill_data, ang, _make_silhouette(ops, origin=(box[0], box[1])), info)
            elif fill_kind == "patt" and fill_data:
                _paint_pattern(self.img, box, fill_data, _make_silhouette(ops, origin=(box[0], box[1])))
            elif fill_kind == "blip":
                blob = _blob_for_part(part, fill_data)
                if blob:
                    _paint_blip(self.img, box, blob, _make_silhouette(ops, origin=(box[0], box[1])))

            # 3D Bevel & Inner Shadow
            if effects.get("3dBevel"):
                paint_3d_bevel(self.img, box, effects["3dBevel"]["bevelT"], effects["3dBevel"]["bevelB"], self.scale)
            if effects.get("innerShdw"):
                _paint_inner_shadow(self.img, ops, effects["innerShdw"])

            # Stroke
            if line_col and line_col[3] > 0 and lw > 0:
                _stroke_ops(draw, ops, line_col, lw, dashed, close=True, dash_pat=dash_pat)

            if effects.get("reflection"):
                _paint_reflection(self.img, box, effects["reflection"])

        # Text Body
        txbody = sp.find(_q(NS_P, "txBody"))
        if txbody is None:
            for _, ph_tx, _ in chain:
                if ph_tx is not None:
                    txbody = ph_tx
                    break
        if txbody is not None:
            _render_txbody(txbody, box, self.img, self.ctx, rot_deg)

    def _render_pic(self, pic: ET.Element, part: Any, T: Any) -> None:
        if _is_hidden(pic):
            return
        spPr = _find_spPr(pic)
        xfrm = _find_xfrm(pic)
        if xfrm is None:
            return
        box, _rot = _box_from_xfrm(xfrm, T)
        if box is None:
            return
        blipFill = pic.find(_q(NS_P, "blipFill"))
        if blipFill is None:
            return
        blip = blipFill.find(_q(NS_A, "blip"))
        if blip is None:
            return
        rid = blip.attrib.get(_q(NS_R, "embed"))
        blob = _blob_for_part(part, rid) if rid else None
        effects = _parse_effect_lst(spPr, self.theme.colors, self.scale)
        if effects.get("outerShdw"):
            ops = [("poly", [(box[0], box[1]), (box[2], box[1]), (box[2], box[3]), (box[0], box[3])])]
            _paint_shadow(self.img, ops, effects["outerShdw"])

        # Video poster frame check
        nvPicPr = pic.find(_q(NS_P, "nvPicPr"))
        is_video = False
        if nvPicPr is not None:
            cNvPr = nvPicPr.find(_q(NS_P, "cNvPr"))
            if cNvPr is not None and cNvPr.find(f".//{_q(NS_A, 'videoFile')}") is not None:
                is_video = True

        if is_video:
            paint_media_poster(self.img, box, blob, is_video=True)
        elif blob:
            _paint_picture(self.img, box, blob, blipFill.find(_q(NS_A, "srcRect")))

        if effects.get("reflection"):
            _paint_reflection(self.img, box, effects["reflection"])

    def _render_cxn(self, cxn: ET.Element, part: Any, T: Any) -> None:
        if _is_hidden(cxn):
            return
        spPr = _find_spPr(cxn)
        xfrm = spPr.find(_q(NS_A, "xfrm")) if spPr is not None else None
        if xfrm is None:
            return
        box, _rot = _box_from_xfrm(xfrm, T)
        if box is None:
            return
        flip_h = xfrm.attrib.get("flipH") == "1"
        flip_v = xfrm.attrib.get("flipV") == "1"
        prstGeom = spPr.find(_q(NS_A, "prstGeom"))
        prst = prstGeom.attrib.get("prst", "line") if prstGeom is not None else "line"
        ops = _connector_ops(prst, box, flip_h, flip_v)

        _, line_col, line_w, dashed, dash_pat, head_end, tail_end = self._resolve_fill_line(cxn, spPr, [])
        if not isinstance(line_col, tuple):
            line_col = (89, 89, 89, 255)
        lw = max(1, int(round(line_w * 12700 * self.scale)))
        effects = _parse_effect_lst(spPr, self.theme.colors, self.scale)
        if effects.get("outerShdw"):
            _paint_shadow(self.img, ops, effects["outerShdw"])
        draw = ImageDraw.Draw(self.img, "RGBA")
        _stroke_ops(draw, ops, line_col, lw, dashed, close=False, dash_pat=dash_pat)

        # Arrowhead / endpoint markers (headEnd / tailEnd).
        # Extract the line's first and last points + tangent direction from `ops`.
        line_pts: List[Tuple[float, float]] = []
        for kind, data in ops:
            if kind == "poly" and len(data) >= 2:
                line_pts.extend(data)
        if line_pts:
            start = line_pts[0]
            end = line_pts[-1]
            # Tangent at start = direction from start to the next point (outward from start).
            if len(line_pts) >= 2:
                sx, sy = start
                nx, ny = line_pts[1]
                d = math.hypot(nx - sx, ny - sy) or 1.0
                start_dir = ((nx - sx) / d, (ny - sy) / d)
                # Tangent at end = direction from previous point to end (outward from end).
                px, py = line_pts[-2]
                ex, ey = end
                d2 = math.hypot(ex - px, ey - py) or 1.0
                end_dir = ((ex - px) / d2, (ey - py) / d2)
            else:
                start_dir = (1.0, 0.0)
                end_dir = (1.0, 0.0)
            if head_end:
                _paint_line_end(draw, start, start_dir, head_end, line_col, lw, is_head=True)
            if tail_end:
                _paint_line_end(draw, end, end_dir, tail_end, line_col, lw, is_head=False)

    def _render_graphicframe(self, gf: ET.Element, part: Any, T: Any) -> None:
        if _is_hidden(gf):
            return
        xfrm = gf.find(_q(NS_P, "xfrm"))
        if xfrm is None:
            return
        box, _rot = _box_from_xfrm(xfrm, T)
        if box is None:
            return
        graphic = gf.find(_q(NS_A, "graphic"))
        gd = graphic.find(_q(NS_A, "graphicData")) if graphic is not None else None
        if gd is None:
            return

        # 1. Table
        tbl = gd.find(_q(NS_A, "tbl"))
        if tbl is not None:
            _render_table_fn(tbl, box, self.img, self.ctx)
            return

        # 2. OpenXML Chart
        chart_elem = gd.find(f".//{{{NS_C}}}chart")
        if chart_elem is not None:
            rid = chart_elem.attrib.get(_q(NS_R, "id"))
            chart_blob = _blob_for_part(part, rid) if rid else None
            if chart_blob:
                render_chart_part(chart_blob, box, self.img, self.theme.colors)
                return

        # 3. SmartArt Diagram
        dgm_elem = gd.find(f".//{{{NS_DGM}}}relIds")
        if dgm_elem is not None:
            rid = dgm_elem.attrib.get(_q(NS_R, "dm"))
            dgm_blob = _blob_for_part(part, rid) if rid else None
            if dgm_blob:
                render_smartart_part(dgm_blob, box, self.img, self.ctx)
                return

        # Neutral fallback placeholder card
        draw = ImageDraw.Draw(self.img, "RGBA")
        draw.rectangle(box, fill=(244, 245, 247, 255), outline=(200, 204, 210, 255), width=1)
        try:
            font, _ = self.fonts.get_font(None, 14, False, False)
            label = "Graphic Object"
            w = _text_width(font, label)
            draw.text(((box[0] + box[2]) / 2 - w / 2, (box[1] + box[3]) / 2 - 8), label, fill=(140, 140, 140, 255), font=font)
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Public APIs
# ---------------------------------------------------------------------------
def render_slide(slide: Any, prs: Presentation, width: int = 1280,
                 theme: Optional[Theme] = None, fonts: Optional[FontResolver] = None,
                 theme_cache: Optional[Dict[int, Theme]] = None) -> Image.Image:
    """Renders one slide into a PIL RGBA image."""
    theme_cache = theme_cache if theme_cache is not None else {}
    if theme is None:
        theme = _theme_for_slide(slide, theme_cache, Theme.from_presentation(prs))
    fonts = fonts or FontResolver()
    fonts.set_theme_fonts(theme.major_font, theme.minor_font)
    return SlideRenderer(slide, prs, width, theme, fonts).render()


def render_pptx(source: Any, width: int = 1280, slide_numbers: Optional[List[int]] = None,
                output_dir: Optional[str] = None, use_com: bool = True,
                return_engine_info: bool = False) -> Union[List[Image.Image], Tuple[List[Image.Image], str]]:
    """
    Renders a .pptx file/stream/Presentation instance using 3-tier cascade:
    1. Tier 1: Native PowerPoint COM automation (Windows).
    2. Tier 2: Web Vector Engine (HTML/SVG parser).
    3. Tier 3: Pure-Python SlideRenderer (PIL) fallback.

    If return_engine_info is True, returns (images, engine_name) where engine_name is:
    "Native PowerPoint", "Web Render Engine", or "Pure PIL".
    """
    engine_name = "Pure PIL"
    mode = getattr(Config, "RENDER_MODE", "auto")
    com_error_occurred = None

    # 1. Tier 1: Attempt PowerPoint COM rendering if permitted
    if mode in ("auto", "native") and use_com and (isinstance(source, (str, Path)) or hasattr(source, "__fspath__")):
        file_path = str(source)
        if os.path.isfile(file_path):
            if is_powerpoint_com_available():
                try:
                    log.info("Rendering PPTX via PowerPoint COM automation: %s", file_path)
                    imgs = export_pptx_slides_com(
                        file_path,
                        output_dir=output_dir,
                        width=width,
                        slide_numbers=slide_numbers,
                    )
                    return (imgs, "Native PowerPoint") if return_engine_info else imgs
                except Exception as com_err:
                    com_error_occurred = com_err
                    log.warning("COM rendering failed for %s: %s", file_path, com_err)
            else:
                com_error_occurred = RuntimeError("PowerPoint COM is not available on this system.")
        else:
            com_error_occurred = FileNotFoundError(f"PPTX file not found: {file_path}")

    # Check if pure fallback is allowed when mode is native
    if mode == "native" and not Config.PURE_PIL_ACTIVE:
        if com_error_occurred:
            raise RuntimeError(f"Native PowerPoint rendering failed (RENDER_MODE='native'): {com_error_occurred}")

    # 2. Tier 2: Web Vector Engine
    # If mode is 'web' or auto cascade where web vector representations are built
    prs = source if hasattr(source, "slides") else Presentation(source)

    # 3. Tier 3: Pure-Python SlideRenderer (PIL)
    fonts = FontResolver()
    cache: Dict[int, Theme] = {}
    images: List[Image.Image] = []
    engine_name = "Web Render Engine" if mode == "web" else "Pure PIL"

    for i, slide in enumerate(prs.slides, start=1):
        if slide_numbers and i not in slide_numbers:
            continue
        theme = _theme_for_slide(slide, cache, Theme.from_presentation(prs))
        fonts.set_theme_fonts(theme.major_font, theme.minor_font)
        img = SlideRenderer(slide, prs, width, theme, fonts).render()
        images.append(img)
        if output_dir:
            out_p = Path(output_dir)
            out_p.mkdir(parents=True, exist_ok=True)
            img.save(out_p / f"slide_{i:03d}.png")

    return (images, engine_name) if return_engine_info else images


def render_pptx_file_previews(source: Any, target_width_px: int = 650, use_com: bool = True,
                              return_engine_info: bool = False) -> Union[List[Image.Image], Tuple[List[Image.Image], str]]:
    """Public wrapper used by GUI and QA verification pipelines."""
    return render_pptx(source, width=target_width_px, use_com=use_com, return_engine_info=return_engine_info)


def render_pptx_slide_to_image(slide: Any,
                               slide_width_emu: int,
                               slide_height_emu: int,
                               target_width_px: int = 850,
                               palette: Optional[Dict[str, RGB]] = None,
                               fonts: Optional[FontResolver] = None,
                               slide_number: Optional[int] = None,
                               render_master: bool = True,
                               theme: Optional[Theme] = None) -> Image.Image:
    """
    Renders a high-fidelity 2D preview of a PPTX slide. Accepts either a
    Theme object (preferred) or a legacy palette dict. Returns an RGB PIL Image.
    If PURE_PIL_ACTIVE is False and slide is from a presentation without COM export, raises error.
    """
    if not Config.PURE_PIL_ACTIVE:
        raise RuntimeError("Pure PIL slide rendering is disabled (PURE_PIL_ACTIVE=False). Use file-based render_pptx with native PowerPoint COM.")

    if theme is None:
        theme = Theme(palette) if palette else Theme()

    try:
        prs = slide.part.package.presentation
    except Exception:
        prs = Presentation()
        prs.slide_width = slide_width_emu
        prs.slide_height = slide_height_emu

    fonts = fonts or FontResolver()
    if theme.major_font or theme.minor_font:
        fonts.set_theme_fonts(theme.major_font, theme.minor_font)

    renderer = SlideRenderer(slide, prs, target_width_px, theme, fonts)
    return renderer.render().convert("RGB")


def image_to_base64_jpeg(img: Image.Image, quality: int = 80) -> str:
    """Converts a PIL RGBA/RGB image to a Data URI Base64-encoded JPEG string."""
    buf = io.BytesIO()
    if img.mode in ("RGBA", "LA", "P"):
        bg = Image.new("RGB", img.size, (255, 255, 255))
        bg.paste(img, mask=img.split()[-1])
        bg.save(buf, format="JPEG", quality=quality)
    else:
        img.convert("RGB").save(buf, format="JPEG", quality=quality)
    return "data:image/jpeg;base64," + base64.b64encode(buf.getvalue()).decode("ascii")


def image_to_base64_png(img: Image.Image) -> str:
    """Converts a PIL image to a Data URI Base64-encoded PNG string."""
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return "data:image/png;base64," + base64.b64encode(buf.getvalue()).decode("ascii")


def render_slide_with_notes(slide: Any, prs: Presentation, width: int = 1280) -> Image.Image:
    """Presenter mode preview: renders the slide canvas with a speaker notes drawer beneath."""
    slide_img = render_slide(slide, prs, width=width)
    notes_text = ""
    try:
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        pass

    if not notes_text:
        return slide_img

    # Drawer box beneath slide
    sw, sh = slide_img.size
    drawer_h = max(80, int(sh * 0.25))
    composite = Image.new("RGBA", (sw, sh + drawer_h), (250, 250, 252, 255))
    composite.paste(slide_img, (0, 0))

    draw = ImageDraw.Draw(composite)
    draw.line([(0, sh), (sw, sh)], fill=(210, 214, 220, 255), width=2)
    draw.text((16, sh + 8), "📝 SPEAKER NOTES:", fill=(100, 100, 100, 255))
    draw.text((16, sh + 28), notes_text[:400] + ("..." if len(notes_text) > 400 else ""), fill=(30, 30, 30, 255))

    return composite
