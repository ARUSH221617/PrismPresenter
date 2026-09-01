"""
web_renderer.py
===============
Vector Web & HTML/SVG Slide Rendering Engine for PPTX Jahat.

Converts python-pptx presentation slides, layout components, custom geometry,
tables, typography, and fills into responsive, standards-compliant HTML5/SVG vector DOM
structures that can be rendered directly in web browsers or converted to raster images.
"""

from __future__ import annotations

import base64
import html
import io
import math
import os
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from pptx_jahat.tools.renderers.color_resolver import RGB, RGBA, ColorResolver
from pptx_jahat.tools.renderers.typography_engine import FontResolver

EMU_PER_INCH = 914400
EMU_PER_PT = 12700


class PPTXWebRenderer:
    """
    Renders PPTX slides into HTML5 vector DOM structures with SVG shapes and CSS styling.
    """

    def __init__(self, slide: Any, prs: Presentation, target_width_px: int = 1280):
        self.slide = slide
        self.prs = prs
        self.target_width_px = target_width_px

        slide_w_emu = getattr(prs, "slide_width", 12192000)
        slide_h_emu = getattr(prs, "slide_height", 6858000)
        self.slide_width_emu = slide_w_emu
        self.slide_height_emu = slide_h_emu

        self.scale = target_width_px / max(slide_w_emu, 1)
        self.target_height_px = int(slide_h_emu * self.scale)

    def emu_to_px(self, emu: int) -> float:
        return emu * self.scale

    def render_html(self) -> str:
        """Renders slide into self-contained HTML/CSS block."""
        bg_css = self._resolve_background_style()

        html_out = [
            f'<div class="pptx-web-slide" style="position:relative;width:{self.target_width_px}px;'
            f'height:{self.target_height_px}px;overflow:hidden;background:{bg_css};'
            f'box-sizing:border-box;font-family:\'Segoe UI\',sans-serif;">'
        ]

        # Render shapes in z-order
        for shape in self.slide.shapes:
            shape_html = self._render_shape(shape)
            if shape_html:
                html_out.append(shape_html)

        html_out.append("</div>")
        return "\n".join(html_out)

    def _resolve_background_style(self) -> str:
        """Determines slide background color or pattern."""
        try:
            bg = self.slide.background
            if bg and bg.fill:
                fill = bg.fill
                if fill.type == 1 and fill.fore_color and fill.fore_color.rgb:  # Solid
                    rgb = fill.fore_color.rgb
                    return f"rgb({rgb[0]},{rgb[1]},{rgb[2]})"
        except Exception:
            pass
        return "#ffffff"

    def _render_shape(self, shape: Any) -> str:
        """Renders a single PPTX shape to HTML/SVG."""
        try:
            left_px = self.emu_to_px(shape.left)
            top_px = self.emu_to_px(shape.top)
            width_px = self.emu_to_px(shape.width)
            height_px = self.emu_to_px(shape.height)
            rot = getattr(shape, "rotation", 0) or 0
        except Exception:
            return ""

        transform_css = f"transform: rotate({rot}deg);" if rot else ""

        # 1. Picture shape
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return self._render_picture(shape, left_px, top_px, width_px, height_px, transform_css)

        # 2. Table shape
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return self._render_table(shape, left_px, top_px, width_px, height_px, transform_css)

        # 3. Standard AutoShape / Textbox / Freeform
        return self._render_autoshape(shape, left_px, top_px, width_px, height_px, transform_css)

    def _render_picture(self, shape: Any, left: float, top: float, w: float, h: float, transform: str) -> str:
        try:
            image_blob = shape.image.blob
            b64_data = base64.b64encode(image_blob).decode("ascii")
            content_type = shape.image.content_type or "image/png"
            src = f"data:{content_type};base64,{b64_data}"
            return (
                f'<img src="{src}" style="position:absolute;left:{left:.1f}px;top:{top:.1f}px;'
                f'width:{w:.1f}px;height:{h:.1f}px;object-fit:cover;{transform}" alt="slide-img" />'
            )
        except Exception:
            return ""

    def _render_table(self, shape: Any, left: float, top: float, w: float, h: float, transform: str) -> str:
        try:
            table = shape.table
            rows_html = []
            for row in table.rows:
                cells_html = []
                for cell in row.cells:
                    text_content = html.escape(cell.text_frame.text.strip())
                    cells_html.append(
                        f'<td style="border:1px solid #272b3c;padding:6px;font-size:12px;'
                        f'color:#1f2937;background:#ffffff;">{text_content}</td>'
                    )
                rows_html.append(f"<tr>{''.join(cells_html)}</tr>")

            return (
                f'<div style="position:absolute;left:{left:.1f}px;top:{top:.1f}px;width:{w:.1f}px;'
                f'height:{h:.1f}px;overflow:hidden;{transform}">'
                f'<table style="width:100%;height:100%;border-collapse:collapse;">{"".join(rows_html)}</table>'
                f'</div>'
            )
        except Exception:
            return ""

    def _render_autoshape(self, shape: Any, left: float, top: float, w: float, h: float, transform: str) -> str:
        # Determine background fill & border
        fill_css = "transparent"
        border_css = "none"

        try:
            if hasattr(shape, "fill") and shape.fill:
                if shape.fill.type == 1 and shape.fill.fore_color and shape.fill.fore_color.rgb:
                    c = shape.fill.fore_color.rgb
                    fill_css = f"rgb({c[0]},{c[1]},{c[2]})"
        except Exception:
            pass

        try:
            if hasattr(shape, "line") and shape.line and shape.line.color and shape.line.color.rgb:
                c = shape.line.color.rgb
                lw = max(1, int(self.emu_to_px(shape.line.width or 12700)))
                border_css = f"{lw}px solid rgb({c[0]},{c[1]},{c[2]})"
        except Exception:
            pass

        # Text content
        text_html = ""
        if hasattr(shape, "has_text_frame") and shape.has_text_frame and shape.text_frame:
            text_html = self._render_text_frame(shape.text_frame, w, h)

        return (
            f'<div style="position:absolute;left:{left:.1f}px;top:{top:.1f}px;width:{w:.1f}px;'
            f'height:{h:.1f}px;background:{fill_css};border:{border_css};box-sizing:border-box;'
            f'display:flex;flex-direction:column;{transform}">'
            f'{text_html}'
            f'</div>'
        )

    def _render_text_frame(self, tf: Any, box_w: float, box_h: float) -> str:
        paragraphs = []
        for p in tf.paragraphs:
            p_text = html.escape(p.text.strip())
            if not p_text:
                continue

            font_size_pt = 14
            font_color_css = "#111827"
            is_bold = False
            is_italic = False

            if p.font:
                try:
                    if p.font.size:
                        font_size_pt = p.font.size.pt
                except Exception:
                    pass
                try:
                    if p.font.bold:
                        is_bold = True
                except Exception:
                    pass
                try:
                    if p.font.italic:
                        is_italic = True
                except Exception:
                    pass
                try:
                    if p.font.color and p.font.color.type == 1 and p.font.color.rgb:
                        c = p.font.color.rgb
                        font_color_css = f"rgb({c[0]},{c[1]},{c[2]})"
                except Exception:
                    pass

            # Detect Persian/Arabic RTL characters
            is_rtl = any("\u0600" <= ch <= "\u06FF" or "\u0750" <= ch <= "\u077F" for ch in p_text)
            dir_attr = 'dir="rtl"' if is_rtl else 'dir="ltr"'
            align_css = "right" if is_rtl else "left"

            font_size_px = max(9.0, font_size_pt * (self.scale * EMU_PER_PT))
            font_weight = "bold" if is_bold else "normal"
            font_style = "italic" if is_italic else "normal"

            paragraphs.append(
                f'<p {dir_attr} style="margin:2px 0;font-size:{font_size_px:.1f}px;color:{font_color_css};'
                f'font-weight:{font_weight};font-style:{font_style};text-align:{align_css};line-height:1.3;">'
                f'{p_text}'
                f'</p>'
            )

        return (
            f'<div style="padding:4px 8px;width:100%;height:100%;box-sizing:border-box;'
            f'overflow:hidden;display:flex;flex-direction:column;justify-content:center;">'
            f'{"".join(paragraphs)}'
            f'</div>'
        )


def render_slide_to_html(slide: Any, prs: Presentation, width: int = 1280) -> str:
    """Renders a single slide into HTML5 vector format."""
    renderer = PPTXWebRenderer(slide, prs, target_width_px=width)
    return renderer.render_html()


def render_pptx_to_html_deck(source: Any, width: int = 1280) -> List[str]:
    """Renders all slides of a presentation into a list of HTML5 vector representations."""
    prs = source if hasattr(source, "slides") else Presentation(source)
    html_slides = []
    for slide in prs.slides:
        html_slides.append(render_slide_to_html(slide, prs, width=width))
    return html_slides
