import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

import docx
from pptx_jahat.config import DATA_DIR, OUTPUT_DIR
from pptx_jahat.tools.pptx_engine import extract_all_templates, inspect_all_templates
from pptx_jahat.tools.pptx_builder import (
    build_pptx_with_agent,
    clone_slide_across_presentations,
    verify_pptx_integrity,
    repair_pptx_package,
    verify_and_auto_heal_pptx,
    _remove_shape,
    _remove_shapes
)
from pptx_jahat.tools.preview import render_pptx_file_previews

def create_sample_template():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Sample Slide 1 (Title)
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    tbox = s1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(2.0))
    p = tbox.text_frame.paragraphs[0]
    p.text = "Corporate Presentation Master Template"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 40, 80)
    
    # Sample Slide 2 (Content with shape)
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    h = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(0.8), Inches(11.3), Inches(1.0))
    h.text_frame.paragraphs[0].text = "Section Overview Header"
    h.text_frame.paragraphs[0].font.size = Pt(28)
    
    b = s2.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(6.0), Inches(4.5))
    bp = b.text_frame.paragraphs[0]
    bp.text = "Sample bullet point description item."
    bp.font.size = Pt(18)
    
    template_path = DATA_DIR / "sample_template.pptx"
    prs.save(str(template_path))
    print(f"Created template: {template_path}")

def create_sample_docx():
    doc = docx.Document()
    doc.add_heading("AI Powered Presentation System", level=1)
    doc.add_paragraph("This is a summary of the next-generation automated presentation generator.")
    
    doc.add_heading("Key Features", level=2)
    doc.add_paragraph("Automatic parsing of hierarchical documents.")
    doc.add_paragraph("Dynamic visual matching with extracted shapes.")
    
    doc.add_heading("Market Metrics", level=2)
    table = doc.add_table(rows=3, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Accuracy"
    table.cell(1, 1).text = "99.4%"
    table.cell(2, 0).text = "Speedup"
    table.cell(2, 1).text = "10x"
    
    docx_path = DATA_DIR / "sample_document.docx"
    doc.save(str(docx_path))
    print(f"Created docx: {docx_path}")

def test_unit_cross_cloning_and_shape_removal():
    print("\n--- Testing Cross-Presentation Slide Cloning & Shape Removal ---")
    source_files = list(DATA_DIR.glob("T*.pptx"))
    if not source_files:
        source_files = [DATA_DIR / "sample_template.pptx"]
        
    prs1 = Presentation(str(source_files[0]))
    target_prs = Presentation()
    target_prs.slide_width = prs1.slide_width
    target_prs.slide_height = prs1.slide_height
    
    # Clone slide 0
    s_cloned = clone_slide_across_presentations(prs1, target_prs, 0)
    orig_shapes_count = len(s_cloned.shapes)
    print(f"Cloned slide with {orig_shapes_count} shapes.")
    
    # Test shape removal
    if orig_shapes_count > 0:
        _remove_shape(s_cloned, 0)
        assert len(s_cloned.shapes) == orig_shapes_count - 1, "Shape removal failed to decrease shape count"
        print(f"Shape removal verified: {orig_shapes_count} -> {len(s_cloned.shapes)}")
        
    test_out = OUTPUT_DIR / "test_cloned_output.pptx"
    target_prs.save(str(test_out))
    assert test_out.exists(), "Output test PPTX does not exist"
    print(f"Verified cross-cloning and shape removal output at {test_out}")

def test_global_template_inspection():
    print("\n--- Testing Global Multi-Template Inspection & Vision Previews ---")
    all_slides = inspect_all_templates(DATA_DIR, include_screenshots=True, screenshot_width=300)
    print(f"Found {len(all_slides)} total candidate slides across templates.")
    assert len(all_slides) > 0, "No candidate slides found across templates"
    
    slides_with_screenshots = [s for s in all_slides if s.get("screenshot_base64")]
    print(f"Generated screenshot base64 payloads for {len(slides_with_screenshots)} slides.")
    assert len(slides_with_screenshots) > 0, "No screenshots were generated"
    assert slides_with_screenshots[0]["screenshot_base64"].startswith("data:image/jpeg;base64,"), "Invalid base64 JPEG format"

def main():
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    create_sample_template()
    create_sample_docx()
    
    test_unit_cross_cloning_and_shape_removal()
    test_global_template_inspection()
    
    print("\n--- Testing Multi-Template Vision Presentation Builder ---")
    out_pptx = build_pptx_with_agent(DATA_DIR / "sample_document.docx", log_callback=print)
    print(f"\nSuccessfully generated presentation: {out_pptx}")
    
    # Verify generated presentation integrity
    is_valid, issues = verify_pptx_integrity(out_pptx)
    print(f"PPTX Integrity check: is_valid={is_valid}, issues={issues}")
    assert is_valid, f"PPTX integrity check failed: {issues}"
    
    # Verify generated presentation can be rendered to preview images
    previews = render_pptx_file_previews(out_pptx, target_width_px=600)
    print(f"Rendered {len(previews)} preview images from output deck.")
    assert len(previews) > 0, "Generated presentation has 0 previewable slides"

    # Run full pure-Python render engine test suite
    print("\n--- Running Pure-Python Slide Rendering Engine Suite ---")
    import unittest
    from tests.test_render_engine import TestRenderEngine
    suite = unittest.TestLoader().loadTestsFromTestCase(TestRenderEngine)
    runner = unittest.TextTestRunner(verbosity=2)
    result = runner.run(suite)
    assert result.wasSuccessful(), "Render engine test suite failed"
    print("ALL TESTS PASSED SUCCESSFULLY!")

if __name__ == "__main__":
    main()
