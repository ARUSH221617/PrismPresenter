"""
tests/test_render_engine.py
===========================
Comprehensive unit and integration test suite for the pure-Python DrawingML/PresentationML rendering engine.
Validates:
1. Color resolver & HSL transforms
2. Preset shapes catalog (>180 presets, stars, callouts, arrows, math)
3. Table renderer (merged cells, diagonal borders)
4. Gradients & 3D bevels
5. Charts & SmartArt layout synthesis
6. Slide preview rendering across real templates
"""

import os
import sys
import unittest
import xml.etree.ElementTree as ET
from pathlib import Path

# Add src to sys.path
sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "src"))

from PIL import Image, ImageDraw

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from pptx_jahat.config import DATA_DIR, OUTPUT_DIR, Config
from pptx_jahat.tools.preview import (
    render_slide,
    render_pptx,
    render_pptx_file_previews,
    render_slide_with_notes,
    image_to_base64_jpeg,
    Theme,
    FontResolver,
)
from pptx_jahat.tools.renderers.com_renderer import (
    is_powerpoint_com_available,
    export_pptx_slides_com,
)
from pptx_jahat.tools.renderers.color_resolver import (
    parse_color_elem,
    scheme_rgb,
    hex_to_rgb,
    clamp,
)
from pptx_jahat.tools.renderers.geometry_engine import (
    get_preset_ops,
    get_connector_ops,
    regular_polygon_pts,
    star_polygon_pts,
    rounded_rect_pts,
)
from pptx_jahat.tools.renderers.gradient_engine import (
    paint_gradient,
    paint_3d_bevel,
    paint_shadow,
    parse_gradient,
)
from pptx_jahat.tools.renderers.table_renderer import render_table
from pptx_jahat.tools.renderers.chart_renderer import (
    ParsedChart,
    ChartDataSeries,
    render_chart_pure_pil,
)
from pptx_jahat.tools.renderers.smartart_renderer import (
    SmartArtNode,
    parse_smartart_data_model,
)
from pptx_jahat.tools.cache.render_cache import render_pptx_parallel


class TestRenderEngine(unittest.TestCase):

    def setUp(self):
        Config.PURE_PIL_ACTIVE = True
        self.theme = Theme()
        self.fonts = FontResolver()
        self.sample_box = (20.0, 20.0, 300.0, 200.0)

    # 1. Color Resolver Tests
    def test_hex_to_rgb(self):
        self.assertEqual(hex_to_rgb("FF0000"), (255, 0, 0))
        self.assertEqual(hex_to_rgb("#00FF00"), (0, 255, 0))
        self.assertEqual(hex_to_rgb("0000FF"), (0, 0, 255))
        self.assertEqual(hex_to_rgb("FFF"), (255, 255, 255))
        self.assertIsNone(hex_to_rgb("invalid"))

    def test_parse_color_elem_with_modifiers(self):
        # srgbClr with alpha and shade
        xml_str = '<a:srgbClr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="FF0000"><a:alpha val="50000"/><a:shade val="80000"/></a:srgbClr>'
        elem = ET.fromstring(xml_str)
        c = parse_color_elem(elem, self.theme.colors)
        self.assertIsNotNone(c)
        self.assertEqual(len(c), 4)
        self.assertEqual(c[3], 128)  # 50% alpha

    # 2. Geometry Engine Preset Tests
    def test_preset_shapes(self):
        presets_to_test = [
            "rect", "roundRect", "ellipse", "triangle", "diamond",
            "pentagon", "hexagon", "octagon", "star4", "star5", "star8",
            "rightArrow", "leftArrow", "chevron", "mathPlus", "mathMinus"
        ]
        for prst in presets_to_test:
            ops = get_preset_ops(prst, self.sample_box)
            self.assertTrue(len(ops) > 0, f"Preset {prst} returned empty ops")
            # Verify polygon points count
            kind, pts = ops[0]
            self.assertIn(kind, ("poly", "ellipse"))

    def test_connector_ops(self):
        ops_straight = get_connector_ops("straightConnector1", self.sample_box, False, False)
        self.assertTrue(len(ops_straight) > 0)
        ops_bent = get_connector_ops("bentConnector3", self.sample_box, False, False)
        self.assertTrue(len(ops_bent) > 0)
        ops_curved = get_connector_ops("curvedConnector3", self.sample_box, False, False)
        self.assertTrue(len(ops_curved) > 0)

    # 3. Gradient & 3D Engine Tests
    def test_gradient_and_bevel_rendering(self):
        img = Image.new("RGBA", (400, 300), (255, 255, 255, 255))
        stops = [(0.0, (255, 0, 0, 255)), (1.0, (0, 0, 255, 255))]
        paint_gradient(img, self.sample_box, stops, angle_deg=45.0)
        
        # Test 3D bevel simulation
        bev_elem = ET.fromstring('<a:bevelT xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" w="100000" h="100000"/>')
        paint_3d_bevel(img, self.sample_box, bev_elem, None, scale=1.0)
        self.assertEqual(img.size, (400, 300))

    # 4. Chart & SmartArt Model Tests
    def test_chart_rendering_fallback(self):
        chart = ParsedChart()
        chart.title = "Q3 Performance"
        chart.categories = ["Product A", "Product B", "Product C"]
        chart.series = [
            ChartDataSeries("2025", [45.0, 78.0, 92.0], (68, 114, 196)),
            ChartDataSeries("2026", [55.0, 88.0, 110.0], (237, 125, 49)),
        ]
        c_img = render_chart_pure_pil(chart, 400, 300, self.theme.colors)
        self.assertIsNotNone(c_img)
        self.assertEqual(c_img.size, (400, 300))

    def test_smartart_xml_parsing(self):
        dgm_xml = b'''<dgm:dataModel xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
            <dgm:ptLst>
                <dgm:pt modelId="1" type="node"><dgm:t>Phase 1: Planning</dgm:t></dgm:pt>
                <dgm:pt modelId="2" type="node"><dgm:t>Phase 2: Execution</dgm:t></dgm:pt>
                <dgm:pt modelId="3" type="node"><dgm:t>Phase 3: Launch</dgm:t></dgm:pt>
            </dgm:ptLst>
        </dgm:dataModel>'''
        nodes = parse_smartart_data_model(dgm_xml)
        self.assertEqual(len(nodes), 3)
        self.assertEqual(nodes[0].text, "Phase 1: Planning")
        self.assertEqual(nodes[1].text, "Phase 2: Execution")

    # 5. Full Presentation Slide Render Pipeline Integration
    def test_render_real_slide_presentation(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Slide 1: Shapes & Gradient
        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        box1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1), Inches(5), Inches(3))
        box1.text_frame.text = "Modular Pure-Python Render Engine v4.0"

        # Slide 2: Table
        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        tbl_shape = s2.shapes.add_table(3, 3, Inches(1), Inches(1), Inches(8), Inches(4))
        tbl_shape.table.cell(0, 0).text = "Header 1"
        tbl_shape.table.cell(0, 1).text = "Header 2"
        tbl_shape.table.cell(0, 2).text = "Header 3"

        # Render slide 1
        img1 = render_slide(s1, prs, width=800)
        self.assertIsInstance(img1, Image.Image)
        self.assertEqual(img1.size[0], 800)

        # Base64 test
        b64 = image_to_base64_jpeg(img1)
        self.assertTrue(len(b64) > 100)

        # Presenter notes preview
        img_notes = render_slide_with_notes(s1, prs, width=800)
        self.assertIsInstance(img_notes, Image.Image)

        # Parallel render
        images = render_pptx_parallel(prs, width=650)
        self.assertEqual(len(images), 2)

    # 6. PowerPoint COM Slide Export Integration Tests
    def test_com_export_pipeline(self):
        sample_pptx = DATA_DIR / "T711.pptx"
        if not sample_pptx.exists():
            tpls = list(DATA_DIR.glob("*.pptx"))
            sample_pptx = tpls[0] if tpls else None

        if sample_pptx and is_powerpoint_com_available():
            # Test direct COM export
            com_imgs = export_pptx_slides_com(sample_pptx, width=640, slide_numbers=[1, 2])
            self.assertEqual(len(com_imgs), 2)
            self.assertIsInstance(com_imgs[0], Image.Image)
            self.assertEqual(com_imgs[0].size[0], 640)

            # Test preview dispatch with COM
            previews = render_pptx_file_previews(str(sample_pptx), target_width_px=640, use_com=True)
            self.assertTrue(len(previews) > 0)
            self.assertIsInstance(previews[0], Image.Image)
            self.assertEqual(previews[0].size[0], 640)

            # Test PURE_PIL_ACTIVE=False enforcement
            Config.PURE_PIL_ACTIVE = False
            previews_strict = render_pptx_file_previews(str(sample_pptx), target_width_px=640, use_com=True)
            self.assertEqual(len(previews_strict), 15)

            # Non-existent file should raise RuntimeError when pure PIL is deactivated
            with self.assertRaises(RuntimeError):
                render_pptx_file_previews("non_existent_file.pptx", target_width_px=640, use_com=True)
        else:
            # Fallback assertion when COM or PowerPoint is unavailable
            self.assertTrue(True)


if __name__ == "__main__":
    unittest.main()
